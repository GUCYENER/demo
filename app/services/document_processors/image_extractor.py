"""
VYRA L1 Support API - Image Extractor Module
=============================================
DOCX, PDF ve PPTX dosyalarından görsel çıkarma.
Görseller DB'de BYTEA olarak saklanır, chunk referansı ile ilişkilendirilir.

Author: VYRA AI Team
Version: 1.0.0
"""

import io
from dataclasses import dataclass
from typing import List, Optional

from app.services.logging_service import log_system_event, log_error


@dataclass
class ExtractedImage:
    """Çıkarılan görsel verisi"""
    image_data: bytes
    image_format: str       # png, jpeg, gif, emf
    width: int = 0
    height: int = 0
    context_heading: str = ""
    context_chunk_index: int = 0
    alt_text: str = ""
    ocr_text: str = ""      # EasyOCR ile çıkarılan metin
    nearby_text: str = ""   # v3.4.5: Görselin yakın çevresindeki metin (chunk eşleştirme için)
    paragraph_index: int = -1    # Görselin orijinal dosyadaki paragraf sırası (0-indexed)
    page_y_position: float = -1  # PDF: sayfa içi Y koordinatı (mm)
    page_number: int = -1        # v3.4.6: Sayfa/slayt no (1-indexed, chunk metadata.page ile uyumlu)


# Singleton OCR reader (lazy load — model yüklemesi ~20sn sadece ilk kez)
_ocr_reader = None


class ImageExtractor:
    """
    DOCX/PDF/PPTX dosyalarından görsel çıkarma modülü.
    
    Kullanım:
        extractor = ImageExtractor()
        images = extractor.extract(file_content, ".docx")
        if images:
            image_ids, saved_images = extractor.save_to_db(images, file_id, cursor)
    """

    # Desteklenen görsel formatları
    SUPPORTED_FORMATS = {"png", "jpeg", "jpg", "gif", "bmp", "tiff", "emf", "wmf"}
    # OCR desteklenen formatlar (EMF/WMF Pillow tarafından açılamayabilir)
    OCR_FORMATS = {"png", "jpeg", "jpg", "gif", "bmp", "tiff"}

    def extract(self, file_content: bytes, file_type: str, skip_ocr: bool = False) -> List[ExtractedImage]:
        """
        Dosya tipine göre görselleri çıkarır.
        
        Args:
            file_content: Dosya binary içeriği
            file_type: Dosya uzantısı (örn: ".docx", ".pdf")
            skip_ocr: True ise OCR atlanır (preview/enhancement için performans)
        
        Returns:
            Çıkarılan görsellerin listesi
        """
        ext = file_type.lower().strip(".")
        
        try:
            if ext in ("docx", "doc"):
                images = self._extract_docx_images(file_content)
            elif ext == "pdf":
                images = self._extract_pdf_images(file_content)
            elif ext in ("pptx", "ppt"):
                images = self._extract_pptx_images(file_content)
            else:
                return []
            
            # v3.4.2: Enhancement preview'da OCR gereksiz — skip_ocr ile atlanabilir
            if images and not skip_ocr:
                self._run_ocr_batch(images)
            
            return images
        except Exception as e:
            log_error(f"Görsel çıkarma hatası ({ext}): {e}", "image_extractor")
            return []

    # ─────────────────────────────────────────
    #  DOCX Görsel Çıkarma
    # ─────────────────────────────────────────

    def _extract_docx_images(self, file_content: bytes) -> List[ExtractedImage]:
        """python-docx ile DOCX'ten görsel çıkar — nearby_text ile bağlam kaydedilir"""
        from docx import Document

        doc = Document(io.BytesIO(file_content))
        images: List[ExtractedImage] = []
        current_heading = "Genel"
        # v3.4.5: Nearby text için son N paragrafı tut
        recent_paragraphs: list = []

        for para_idx, para in enumerate(doc.paragraphs):
            style_name = para.style.name if para.style else ""
            para_text = para.text.strip()

            # Heading takibi
            if style_name.startswith("Heading") and para_text:
                current_heading = para_text

            # Son 5 paragrafı nearby_text için tut
            if para_text:
                recent_paragraphs.append(para_text)
                if len(recent_paragraphs) > 5:
                    recent_paragraphs.pop(0)

            # v3.4.7: InlineShape VE Anchor (floating) görselleri kontrol et
            WP_NS = '{http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing}'
            A_NS = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
            R_NS = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
            
            for run in para.runs:
                # Issue 1: Hem inline hem anchor (floating) görselleri tara
                drawing_elements = []
                
                # InlineShape (gömülü resimler)
                inlines = run._element.findall(f'.//{WP_NS}inline')
                drawing_elements.extend(inlines)
                
                # Anchor (floating, sağa/sola hizalanmış resimler)
                anchors = run._element.findall(f'.//{WP_NS}anchor')
                drawing_elements.extend(anchors)
                
                for drawing_el in drawing_elements:
                    blip = drawing_el.find(f'.//{A_NS}blip')
                    if blip is not None:
                        embed_id = blip.get(f'{R_NS}embed')
                        if embed_id:
                            nearby = " ".join(recent_paragraphs)[:300]
                            img = self._get_docx_image_from_rel(
                                doc, embed_id, current_heading, para_idx, nearby
                            )
                            if img:
                                images.append(img)

        log_system_event("INFO", f"DOCX'ten {len(images)} görsel çıkarıldı (inline+anchor)", "image_extractor")
        return images

    def _get_docx_image_from_rel(self, doc, embed_id: str, heading: str, para_idx: int = -1, nearby_text: str = "") -> Optional[ExtractedImage]:
        """DOCX relationship'ten görsel verisini alır"""
        try:
            rel = doc.part.rels.get(embed_id)
            if rel is None:
                return None

            image_part = rel.target_part
            image_data = image_part.blob
            content_type = image_part.content_type or ""

            # Format belirle
            fmt = self._format_from_content_type(content_type)
            if not fmt:
                return None

            width, height = self._get_image_dimensions(image_data, fmt)

            # Çok küçük görselleri atla (ikon, decoration vb.)
            if width < 50 or height < 50:
                return None

            return ExtractedImage(
                image_data=image_data,
                image_format=fmt,
                width=width,
                height=height,
                context_heading=heading[:500],
                context_chunk_index=-1,  # v3.4.5: Metin bazlı eşleştirme kullanılacak
                alt_text=f"Görsel - {heading[:100]}",
                nearby_text=nearby_text,
                paragraph_index=para_idx
            )
        except Exception as e:
            log_error(f"DOCX görsel okuma hatası: {e}", "image_extractor")
            return None

    # ─────────────────────────────────────────
    #  PDF Görsel Çıkarma
    # ─────────────────────────────────────────

    def _extract_pdf_images(self, file_content: bytes) -> List[ExtractedImage]:
        """PyMuPDF (fitz) ile PDF'ten görsel çıkar — sayfa heading'leri ve Y pozisyonu kaydeder"""
        import fitz
        import re

        doc = fitz.open(stream=file_content, filetype="pdf")
        images: List[ExtractedImage] = []
        
        # Toplam metin satır sayacı (yapay paragraph_index olarak kullanılacak)
        global_line_counter = 0

        for page_num in range(len(doc)):
            page = doc[page_num]
            image_list = page.get_images(full=True)
            
            if not image_list:
                # Sayfadaki satır sayısını sayacıya ekle (görsel olmasa bile)
                page_text = page.get_text("text") or ""
                global_line_counter += len([l for l in page_text.split('\n') if l.strip()])
                continue
            
            # Sayfadaki metin bloklarını pozisyonlarıyla birlikte al
            page_text = page.get_text("text") or ""
            page_lines = [l.strip() for l in page_text.split('\n') if l.strip()]
            
            # Heading'leri satır pozisyonu (metin bloğu dict) ile topla
            headings_on_page = []
            try:
                text_blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)
                for block in text_blocks.get("blocks", []):
                    if block.get("type") != 0:  # Sadece text blokları
                        continue
                    for line_info in block.get("lines", []):
                        spans = line_info.get("spans", [])
                        if not spans:
                            continue
                        line_text = "".join(s.get("text", "") for s in spans).strip()
                        if not line_text or len(line_text) > 120 or len(line_text) < 3:
                            continue
                        # v3.4.7: Heading tespiti iyileştirildi
                        avg_size = sum(s.get("size", 11) for s in spans) / len(spans)
                        is_bold = any(s.get("flags", 0) & 2**4 for s in spans)
                        y_pos = line_info["bbox"][1]  # Top Y
                        
                        # Font size >= 12, bold, veya font size >= 14 (başlık formatı)
                        if avg_size >= 12 or is_bold or avg_size >= 14:
                            # Heading olabilecek metin kontrolü:
                            # - Sayfa numarası değil, cümle ortası değil
                            # - Numara ile başlayabilir (1.2 Giriş gibi)
                            is_page_num = re.match(r'^\d+$', line_text.strip())
                            ends_with_sentence = line_text.rstrip().endswith(('.', ',', ';'))
                            if (not is_page_num and
                                not ends_with_sentence and
                                re.match(r'^[\d\.\s]*[A-ZÇĞİÖŞÜa-zçğıöşü]', line_text)):
                                headings_on_page.append((y_pos, line_text))
            except Exception:
                # Fallback: basit heading tespiti
                for line in page_lines:
                    if (len(line) < 80 and line[0].isupper() and 
                        not line.endswith('.') and
                        re.match(r'^[\d\.\s]*[A-ZÇĞİÖŞÜa-zçğıöşü]', line)):
                        headings_on_page.append((0, line))
            
            # Görsellerin sayfa üzerindeki bounding box'larını al
            for img_idx, img_info in enumerate(image_list):
                xref = img_info[0]
                try:
                    base_image = doc.extract_image(xref)
                    if not base_image or not base_image.get("image"):
                        continue

                    image_data = base_image["image"]
                    fmt = base_image.get("ext", "png").lower()

                    if fmt not in self.SUPPORTED_FORMATS:
                        continue

                    # Çok küçük görselleri atla (ikon, decoration vb.)
                    width = base_image.get("width", 0)
                    height = base_image.get("height", 0)
                    if width < 50 or height < 50:
                        continue
                    
                    # Görselin sayfa üzerindeki Y pozisyonunu bul
                    y_pos = 0.0
                    try:
                        img_rects = page.get_image_rects(xref)
                        if img_rects:
                            y_pos = img_rects[0].y0  # İlk rect'in top Y
                    except Exception:
                        y_pos = float(img_idx * 100)  # Fallback: sıra bazlı

                    # Bu görselin hemen öncesindeki heading'i bul
                    best_heading = f"Sayfa {page_num + 1}"
                    for h_y, h_text in sorted(headings_on_page, key=lambda x: x[0], reverse=True):
                        if h_y < y_pos:
                            best_heading = h_text
                            break
                    
                    # v3.4.5: Nearby text — görselin Y pozisyonuna yakın ±5 satırlık metin
                    nearby = ""
                    if page_lines and y_pos > 0:
                        page_height = page.rect.height
                        approx_line = int((y_pos / page_height) * len(page_lines)) if page_height > 0 else 0
                        start_l = max(0, approx_line - 5)
                        end_l = min(len(page_lines), approx_line + 5)
                        nearby = " ".join(page_lines[start_l:end_l])[:300]
                    elif page_lines:
                        nearby = " ".join(page_lines[:10])[:300]
                    
                    # Yapay paragraph_index
                    lines_before_img = sum(1 for l in page_lines if l.strip())
                    if y_pos > 0:
                        page_height = page.rect.height
                        lines_before_img = int((y_pos / page_height) * len(page_lines)) if page_height > 0 else 0
                    
                    para_idx = global_line_counter + lines_before_img

                    images.append(ExtractedImage(
                        image_data=image_data,
                        image_format=fmt if fmt != "jpg" else "jpeg",
                        width=width,
                        height=height,
                        context_heading=best_heading,
                        context_chunk_index=-1,  # v3.4.5: Sayfa no ≠ chunk index, fallback kaldırıldı
                        alt_text=f"PDF Sayfa {page_num + 1} - Görsel {img_idx + 1}",
                        nearby_text=nearby,
                        paragraph_index=para_idx,
                        page_y_position=y_pos,
                        page_number=page_num + 1  # 1-indexed (chunk metadata.page ile uyumlu)
                    ))
                except Exception as e:
                    log_error(f"PDF görsel çıkarma hatası (sayfa {page_num}): {e}", "image_extractor")

            # Bu sayfanın satır sayısını toplam sayaca ekle
            global_line_counter += len(page_lines)

        doc.close()
        log_system_event("INFO", f"PDF'ten {len(images)} görsel çıkarıldı", "image_extractor")
        return images

    # ─────────────────────────────────────────
    #  PPTX Görsel Çıkarma
    # ─────────────────────────────────────────

    def _extract_pptx_images(self, file_content: bytes) -> List[ExtractedImage]:
        """python-pptx ile PPTX'ten görsel çıkar"""
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(io.BytesIO(file_content))
        images: List[ExtractedImage] = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_title = ""
            # Slayt başlığını bul
            if slide.shapes.title:
                slide_title = slide.shapes.title.text or f"Slayt {slide_idx + 1}"

            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        image_data = image.blob
                        content_type = image.content_type or ""

                        fmt = self._format_from_content_type(content_type)
                        if not fmt:
                            continue

                        width, height = self._get_image_dimensions(image_data, fmt)

                        # Çok küçük görselleri atla
                        if width < 50 or height < 50:
                            continue

                        # v3.4.5: Slayttaki tüm metin shape'lerini nearby_text olarak topla
                        slide_texts = []
                        for s in slide.shapes:
                            if s.has_text_frame:
                                for p in s.text_frame.paragraphs:
                                    t = p.text.strip()
                                    if t:
                                        slide_texts.append(t)
                        slide_nearby = " ".join(slide_texts)[:300]
                        
                        images.append(ExtractedImage(
                            image_data=image_data,
                            image_format=fmt,
                            width=width,
                            height=height,
                            context_heading=(slide_title or f"Slayt {slide_idx + 1}")[:500],
                            context_chunk_index=-1,  # v3.4.5: Metin bazlı eşleştirme kullanılacak
                            alt_text=f"Slayt {slide_idx + 1} - {slide_title[:100]}",
                            nearby_text=slide_nearby,
                            page_number=slide_idx + 1  # Slayt no (1-indexed)
                        ))
                    except Exception as e:
                        log_error(f"PPTX görsel çıkarma hatası (slayt {slide_idx}): {e}", "image_extractor")

        log_system_event("INFO", f"PPTX'ten {len(images)} görsel çıkarıldı", "image_extractor")
        return images

    # ─────────────────────────────────────────
    #  Yardımcı Fonksiyonlar
    # ─────────────────────────────────────────

    def _format_from_content_type(self, content_type: str) -> Optional[str]:
        """Content-Type'tan görsel formatını belirle"""
        ct_map = {
            "image/png": "png",
            "image/jpeg": "jpeg",
            "image/jpg": "jpeg",
            "image/gif": "gif",
            "image/bmp": "bmp",
            "image/tiff": "tiff",
            "image/x-emf": "emf",
            "image/x-wmf": "wmf",
        }
        return ct_map.get(content_type.lower(), None)

    def _get_image_dimensions(self, image_data: bytes, fmt: str) -> tuple:
        """Görsel boyutlarını Pillow ile hesapla. EMF/WMF için fallback boyut döner."""
        # v3.4.7 Issue 2: EMF/WMF Pillow tarafından açılamayabilir — makul varsayılan döndür
        if fmt and fmt.lower() in ('emf', 'wmf'):
            # Vektörel format — Pillow genelde başarısız olur
            # Varsayılan boyut: 300x200 (ortalama diagram boyutu)
            try:
                from PIL import Image
                img = Image.open(io.BytesIO(image_data))
                w, h = img.size
                if w > 0 and h > 0:
                    return (w, h)
            except Exception:
                pass
            return (300, 200)
        
        try:
            from PIL import Image
            img = Image.open(io.BytesIO(image_data))
            return img.size  # (width, height)
        except Exception as e:
            log_error(f"Image boyut okuma hatası ({fmt}): {e}", "image_extractor")
            return (0, 0)

    # ─────────────────────────────────────────
    #  OCR — EasyOCR ile Metin Çıkarma
    # ─────────────────────────────────────────

    def _get_ocr_reader(self):
        """Lazy singleton EasyOCR reader (ilk çağrıda model yüklenir)"""
        global _ocr_reader
        if _ocr_reader is None:
            try:
                import easyocr
                _ocr_reader = easyocr.Reader(['tr', 'en'], gpu=False, verbose=False)
                log_system_event("INFO", "EasyOCR reader yüklendi (tr+en)", "image_extractor")
            except ImportError:
                log_error("EasyOCR yüklü değil. pip install easyocr", "image_extractor")
                return None
            except Exception as e:
                log_error(f"EasyOCR yükleme hatası: {e}", "image_extractor")
                return None
        return _ocr_reader

    def _run_ocr_single(self, image_data: bytes, fmt: str) -> str:
        """
        Tek görsel için OCR çalıştır.
        
        Returns:
            Çıkarılan metin (boş string hata durumunda)
        """
        if fmt not in self.OCR_FORMATS:
            return ""
        
        reader = self._get_ocr_reader()
        if reader is None:
            return ""
        
        try:
            import numpy as np
            from PIL import Image
            
            img = Image.open(io.BytesIO(image_data))
            # RGB'ye çevir (RGBA, P vb. modlar EasyOCR'da sorun yaratabilir)
            if img.mode != 'RGB':
                img = img.convert('RGB')
            
            img_array = np.array(img)
            results = reader.readtext(img_array, detail=0, paragraph=True)
            
            text = "\n".join(results).strip()
            return text
        except Exception as e:
            log_error(f"OCR hatası: {e}", "image_extractor")
            return ""

    def _run_ocr_batch(self, images: List[ExtractedImage]):
        """
        Tüm görseller için OCR çalıştır ve ocr_text alanını doldur.
        ThreadPoolExecutor ile paralel çalışır (max 4 thread).
        """
        from concurrent.futures import ThreadPoolExecutor, as_completed
        
        if not images:
            return
        
        # İlk çağrıda model yüklenmesini sağla (thread-safe init)
        if self._get_ocr_reader() is None:
            log_system_event("WARNING", "OCR reader yüklenemedi, batch atlanıyor", "image_extractor")
            return
        
        max_workers = min(4, len(images))
        
        def _ocr_task(img):
            try:
                return self._run_ocr_single(img.image_data, img.image_format)
            except Exception as e:
                log_error(f"OCR task hatası ({img.image_format}): {e}", "image_extractor")
                return ""
        
        ocr_count = 0
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            future_to_img = {executor.submit(_ocr_task, img): img for img in images}
            for future in as_completed(future_to_img):
                img = future_to_img[future]
                try:
                    text = future.result()
                    if text:
                        img.ocr_text = text
                        ocr_count += 1
                except Exception as e:
                    log_error(f"OCR future hatası ({img.image_format}): {e}", "image_extractor")
        
        if ocr_count > 0:
            log_system_event("INFO", f"{ocr_count}/{len(images)} görselden OCR metin çıkarıldı (paralel)", "image_extractor")
        elif len(images) > 0:
            log_system_event("WARNING", f"0/{len(images)} görselden OCR metin çıkarılamadı", "image_extractor")

    def save_to_db(self, images: List[ExtractedImage], file_id: int, cursor) -> tuple:
        """
        Görselleri document_images tablosuna kaydet.
        
        Args:
            images: Çıkarılan görseller
            file_id: Dosya ID
            cursor: DB cursor (dışarıdan geçirilir, commit yapılmaz)
        
        Returns:
            Tuple[List[int], List[ExtractedImage]]: (image_ids, başarıyla kaydedilen görseller)
            v3.4.6: Hata alan görseller hariç tutulur — zip hizalama güvenliği.
        """
        if not images:
            return [], []

        image_ids = []
        successful_images = []
        for idx, img in enumerate(images):
            try:
                cursor.execute(
                    """
                    INSERT INTO document_images 
                    (file_id, image_index, image_data, image_format, 
                     width_px, height_px, file_size_bytes, 
                     context_heading, context_chunk_index, alt_text, ocr_text)
                    VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
                    RETURNING id
                    """,
                    (
                        file_id, idx, img.image_data, img.image_format,
                        img.width, img.height, len(img.image_data),
                        img.context_heading, img.context_chunk_index, img.alt_text,
                        img.ocr_text
                    )
                )
                row = cursor.fetchone()
                if row:
                    image_ids.append(row["id"])
                    successful_images.append(img)
            except Exception as e:
                log_error(f"Görsel DB kayıt hatası (index {idx}): {e}", "image_extractor")

        log_system_event("INFO", f"Dosya {file_id} için {len(image_ids)} görsel kaydedildi", "image_extractor")
        return image_ids, successful_images
