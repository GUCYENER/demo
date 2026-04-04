"""
PowerPoint Document Processor
Uses python-pptx for PPTX text extraction
Enhanced with slide title extraction, speaker notes, and semantic chunking (2024 Best Practices)
"""

import logging
from pathlib import Path
from typing import BinaryIO, List

from .base import BaseDocumentProcessor

logger = logging.getLogger("vyra")


class PPTXProcessor(BaseDocumentProcessor):
    """PowerPoint dosyalarını işleyen processor (Enhanced)"""
    
    SUPPORTED_EXTENSIONS = ['.pptx', '.ppt']
    PROCESSOR_NAME = "PPTXProcessor"
    
    def extract_text(self, file_path: Path) -> str:
        """PPTX dosyasından metin çıkarır"""
        try:
            from pptx import Presentation
            
            prs = Presentation(str(file_path))
            return self._extract_from_presentation(prs)
            
        except ImportError:
            raise ImportError("python-pptx kütüphanesi yüklü değil. 'pip install python-pptx' komutunu çalıştırın.")
        except Exception as e:
            raise RuntimeError(f"PPTX işleme hatası: {str(e)}")
    
    def extract_text_from_bytes(self, file_obj: BinaryIO, file_name: str) -> str:
        """BytesIO'dan PPTX metni çıkarır"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_obj)
            return self._extract_from_presentation(prs)
            
        except ImportError:
            raise ImportError("python-pptx kütüphanesi yüklü değil.")
        except Exception as e:
            raise RuntimeError(f"PPTX işleme hatası: {str(e)}")
    
    def _extract_from_presentation(self, prs) -> str:
        """Presentation objesinden metin çıkarır (ortak mantık)"""
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"[Slayt {slide_num}]"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                
                # Tablo içeriğini çıkar
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            slide_text.append(row_text)
            
            if len(slide_text) > 1:
                text_parts.append("\n".join(slide_text))
        
        result = "\n\n".join(text_parts)
        
        # 🇹🇷 Türkçe karakter düzeltme + NUL temizleme
        return self._fix_turkish_chars(result)
    
    # =========================================================================
    # SLIDE TITLE & SEMANTIC CHUNKING (2024 Best Practices)
    # =========================================================================
    
    def _get_slide_title(self, slide) -> str:
        """
        Slayt başlığını çıkarır.
        Title placeholder veya ilk büyük text shape'i kullanır.
        """
        # Önce title placeholder ara
        if slide.shapes.title:
            title_text = slide.shapes.title.text.strip()
            if title_text:
                return title_text
        
        # Title placeholder yoksa, en büyük font'lu kısa metni bul
        best_title = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                # Kısa ve muhtemelen başlık olan metinler
                if len(text) < 100 and "\n" not in text:
                    if not best_title or len(text) > len(best_title):
                        best_title = text
        
        return best_title
    
    def _get_speaker_notes(self, slide) -> str:
        """Slayt konuşmacı notlarını çıkarır."""
        try:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_text = notes_slide.notes_text_frame.text.strip()
                return notes_text
        except Exception as e:
            logger.debug("[PPTXProcessor] Speaker notes okuma hatası", exc_info=True)
        return ""
    
    def _extract_slide_content(self, slide) -> dict:
        """
        Tek bir slayttan zengin içerik çıkarır.
        
        Returns:
            dict: {"title": "...", "content": "...", "notes": "...", "has_table": bool}
        """
        title = self._get_slide_title(slide)
        content_parts = []
        has_table = False
        
        for shape in slide.shapes:
            # Tablo kontrolü
            if shape.has_table:
                has_table = True
                table = shape.table
                table_rows = []
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        table_rows.append(row_text)
                if table_rows:
                    content_parts.append("\n".join(table_rows))
            
            # Normal text content
            elif hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                # Title'ı tekrar ekleme
                if text != title:
                    content_parts.append(text)
        
        notes = self._get_speaker_notes(slide)
        
        return {
            "title": title,
            "content": "\n".join(content_parts),
            "notes": notes,
            "has_table": has_table
        }
    
    def extract_chunks(self, file_obj: BinaryIO, file_name: str) -> List[dict]:
        """
        PPTX dosyasından zengin metadata ile chunk'lar çıkarır.
        PDF/DOCX processor ile tutarlı format.
        
        Returns:
            List[dict]: [{"text": "...", "metadata": {"heading": "...", "slide": N, "type": "..."}}, ...]
        """
        try:
            from pptx import Presentation
            
            # File position'ı başa al
            if hasattr(file_obj, 'seek'):
                file_obj.seek(0)
            
            prs = Presentation(file_obj)
            
        except Exception as e:
            # Fallback - boş liste dön, base.py standart chunking yapacak
            logger.warning("[PPTXProcessor] extract_chunks hatası", exc_info=True)
            return []
        
        chunks = []
        chunk_index = 0
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_data = self._extract_slide_content(slide)
            
            title = slide_data["title"]
            content = slide_data["content"]
            notes = slide_data["notes"]
            has_table = slide_data["has_table"]
            
            # Content type belirleme
            content_type = "table" if has_table else "slide"
            
            # Ana içerik chunk'ı
            if content and len(content.strip()) >= 30:
                # Çok büyük içeriği böl
                if len(content) > 800:
                    parts = self._split_content(content, max_size=800)
                    for part in parts:
                        if len(part.strip()) >= 30:
                            # v3.4.1: Alt-chunk'ta sub-heading tespiti
                            effective_heading = title
                            sub_head = self._detect_sub_heading_text(part)
                            if sub_head and sub_head != title:
                                effective_heading = sub_head
                            
                            chunks.append({
                                "text": part.strip(),
                                "metadata": {
                                    "type": content_type,
                                    "heading": effective_heading,
                                    "slide": slide_num,
                                    "file_type": "pptx",
                                    "chunk_index": chunk_index,
                                    "source": file_name
                                }
                            })
                            chunk_index += 1
                else:
                    chunks.append({
                        "text": content.strip(),
                        "metadata": {
                            "type": content_type,
                            "heading": title,
                            "slide": slide_num,
                            "file_type": "pptx",
                            "chunk_index": chunk_index,
                            "source": file_name
                        }
                    })
                    chunk_index += 1
            
            # Speaker notes chunk'ı (ayrı chunk olarak)
            if notes and len(notes.strip()) >= 50:
                chunks.append({
                    "text": notes.strip(),
                    "metadata": {
                        "type": "speaker_notes",
                        "heading": f"{title} - Notlar" if title else "Konuşmacı Notları",
                        "slide": slide_num,
                        "file_type": "pptx",
                        "chunk_index": chunk_index,
                        "source": file_name
                    }
                })
                chunk_index += 1
        
        return chunks
    
    def _detect_sub_heading_text(self, text: str):
        """
        v3.4.1: Alt-chunk içindeki heading satırını tespit eder.
        İlk 3 satıra bakarak Title Case veya uppercase heading arar.
        """
        if not text or len(text.strip()) < 10:
            return None
        
        lines = text.strip().split('\n')
        for line in lines[:3]:
            stripped = line.strip()
            if not stripped or len(stripped) < 3 or len(stripped) > 80:
                continue
            if stripped.endswith('.'):
                continue
            # Tüm büyük harf → heading
            if stripped.isupper() and len(stripped) < 60:
                return stripped
            # Title Case kontrolü
            words = stripped.split()
            if 2 <= len(words) <= 10:
                tc_count = sum(1 for w in words if len(w) > 1 and w[0].isupper())
                if tc_count / len(words) >= 0.7:
                    return stripped
        return None
    
    def _split_content(self, text: str, max_size: int = 800) -> List[str]:
        """Büyük içeriği paragraf sınırlarında böler."""
        if len(text) <= max_size:
            return [text]
        
        chunks = []
        paragraphs = text.split('\n\n')
        
        current_chunk = ""
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            
            if len(current_chunk) + len(para) + 2 > max_size:
                if current_chunk:
                    chunks.append(current_chunk)
                current_chunk = para
            else:
                current_chunk = (current_chunk + "\n\n" + para).strip()
        
        if current_chunk:
            chunks.append(current_chunk)
        
        return chunks
    
    def get_metadata(self, file_path: Path = None, file_name: str = None) -> dict:
        """PPTX metadata'sını çıkarır"""
        base_meta = {"processor": self.PROCESSOR_NAME}
        
        if file_path:
            try:
                from pptx import Presentation
                prs = Presentation(str(file_path))
                
                base_meta.update({
                    "slide_count": len(prs.slides),
                    "file_name": file_path.name,
                })
            except Exception as e:
                logger.warning("[PPTXProcessor] Metadata okuma hatası", exc_info=True)
                base_meta["file_name"] = file_path.name
        else:
            base_meta["file_name"] = file_name
        
        return base_meta

