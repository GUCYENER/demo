"""
VYRA L1 Support API - Document Enhancement Service
===================================================
CatBoost chunk kalite analizi + LLM ile doküman iyileştirme.

Akış:
1. Maturity ihlallerini al
2. Doküman içeriğini bölümlere ayır
3. CatBoost ile her bölümün RAG kalitesini tahmin et
4. LLM ile düşük kaliteli bölümleri iyileştir
5. İyileştirilmiş DOCX oluştur

Author: VYRA AI Team
Version: 1.0.0 (v2.35.0)
"""

import io
import os
import re
import json
import time
import tempfile
import uuid
from typing import Dict, Any, List, BinaryIO, Optional
from dataclasses import dataclass, field

from app.services.logging_service import log_system_event, log_error, log_warning
from app.services.content_integrity_validator import get_integrity_validator


# ============================================
# Enhancement Data Classes
# ============================================

@dataclass
class EnhancedSection:
    """Bir bölümün orijinal ve iyileştirilmiş hali"""
    section_index: int
    heading: str
    original_text: str
    enhanced_text: str
    change_type: str          # "heading_added", "content_restructured", "table_fixed", "encoding_fixed", "no_change", "llm_error", "integrity_failed"
    explanation: str          # İyileştirme açıklaması
    priority: float           # CatBoost priority skoru (0-1)
    violations: List[str] = field(default_factory=list)
    integrity_score: float = 1.0    # Bütünlük doğrulama skoru (0-1)
    integrity_issues: List[str] = field(default_factory=list)  # Bütünlük sorunları


@dataclass
class EnhancementResult:
    """Tüm iyileştirme sonucu"""
    file_name: str
    file_type: str
    total_sections: int
    enhanced_count: int
    sections: List[EnhancedSection] = field(default_factory=list)
    catboost_summary: Dict[str, Any] = field(default_factory=dict)
    session_id: str = ""
    enhanced_docx_path: str = ""
    error: Optional[str] = None


# ============================================
# Geçici dosya deposu (session bazlı)
# ============================================
_enhanced_files: Dict[str, str] = {}  # session_id → temp file path


# ============================================
# Document Enhancer Service
# ============================================

class DocumentEnhancer:
    """
    CatBoost + LLM tabanlı doküman iyileştirme servisi.
    
    Akış:
    1. _extract_sections(): Dökümanı bölümlere ayır
    2. _catboost_prioritize(): CatBoost ile kalite tahmini
    3. _llm_enhance(): LLM ile iyileştirme
    4. generate_enhanced_docx(): İndirilecek DOCX oluştur
    """
    
    def __init__(self):
        self._catboost_service = None
        self._feature_extractor = None
    
    def _get_catboost(self):
        """CatBoost servisini lazy load et"""
        if self._catboost_service is None:
            try:
                from app.services.catboost_service import get_catboost_service
                self._catboost_service = get_catboost_service()
            except Exception as e:
                log_warning(f"CatBoost yüklenemedi: {e}", "enhancer")
        return self._catboost_service
    
    def _get_feature_extractor(self):
        """Feature extractor'ı lazy load et"""
        if self._feature_extractor is None:
            try:
                from app.services.feature_extractor import get_feature_extractor
                self._feature_extractor = get_feature_extractor()
            except Exception as e:
                log_warning(f"FeatureExtractor yüklenemedi: {e}", "enhancer")
        return self._feature_extractor
    
    # ─────────────────────────────────────────
    #  ANA PİPELINE
    # ─────────────────────────────────────────
    
    def analyze_and_enhance(
        self,
        file_content: bytes,
        file_name: str,
        maturity_result: Dict[str, Any],
        progress_callback=None
    ) -> EnhancementResult:
        """
        Ana pipeline: Analiz → CatBoost → LLM → DOCX
        
        Args:
            file_content: Dosya binary içeriği
            file_name: Dosya adı
            maturity_result: Maturity analyzer sonucu (violations, categories vb.)
        
        Returns:
            EnhancementResult
        """
        file_type = maturity_result.get("file_type", "").upper()
        violations = maturity_result.get("violations", [])
        
        log_system_event("INFO", f"Enhancement başlatıldı: {file_name} ({file_type})", "enhancer")
        start_time = time.time()
        
        try:
            # Adım 1: Bölümlere ayır
            sections = self._extract_sections(file_content, file_name, file_type)
            
            if not sections:
                return EnhancementResult(
                    file_name=file_name,
                    file_type=file_type,
                    total_sections=0,
                    enhanced_count=0,
                    error="Doküman içeriği çıkarılamadı."
                )
            
            log_system_event("INFO", f"{len(sections)} bölüm çıkarıldı", "enhancer")
            
            # Adım 2: CatBoost ile priority analizi
            catboost_analysis = self._catboost_prioritize(sections, file_type, violations)
            
            # Adım 3: LLM ile iyileştirme
            enhanced_sections = self._llm_enhance(sections, violations, file_type, catboost_analysis, progress_callback=progress_callback)
            
            # Adım 4: DOCX oluştur (orijinal formatı koruyarak)
            session_id = str(uuid.uuid4())[:8]
            docx_path = self._generate_enhanced_docx(
                enhanced_sections, file_name, session_id,
                original_content=file_content, file_type=file_type
            )
            
            # Sonucu derle
            enhanced_count = sum(1 for s in enhanced_sections if s.change_type not in ("no_change", "llm_error", "integrity_failed"))
            error_count = sum(1 for s in enhanced_sections if s.change_type in ("llm_error", "integrity_failed"))
            
            result = EnhancementResult(
                file_name=file_name,
                file_type=file_type,
                total_sections=len(sections),
                enhanced_count=enhanced_count,
                sections=enhanced_sections,
                catboost_summary=catboost_analysis.get("summary", {}),
                session_id=session_id,
                enhanced_docx_path=docx_path
            )
            
            elapsed = round(time.time() - start_time, 2)
            log_system_event("INFO", f"Enhancement tamamlandı: {enhanced_count}/{len(sections)} bölüm, {elapsed}s", "enhancer")
            
            return result
            
        except Exception as e:
            log_error(f"Enhancement hatası: {e}", "enhancer")
            return EnhancementResult(
                file_name=file_name,
                file_type=file_type,
                total_sections=0,
                enhanced_count=0,
                error=str(e)
            )
    
    # ─────────────────────────────────────────
    #  ADIM 1: Bölümlere Ayırma
    # ─────────────────────────────────────────
    
    def _extract_sections(
        self,
        file_content: bytes,
        file_name: str,
        file_type: str
    ) -> List[Dict[str, Any]]:
        """Dokümanı bölümlere ayırır (heading bazlı veya sayfa bazlı)"""
        
        file_obj = io.BytesIO(file_content)
        sections = []
        
        try:
            if file_type == "PDF":
                sections = self._extract_pdf_sections(file_obj)
            elif file_type == "DOCX":
                sections = self._extract_docx_sections(file_obj)
            elif file_type == "XLSX":
                sections = self._extract_xlsx_sections(file_obj)
            elif file_type == "PPTX":
                sections = self._extract_pptx_sections(file_obj)
            elif file_type == "CSV":
                sections = self._extract_csv_sections(file_obj)
            elif file_type == "TXT":
                sections = self._extract_txt_sections(file_obj)
            else:
                # Bilinmeyen format — düz metin olarak dene
                file_obj.seek(0)
                text = file_obj.read().decode("utf-8", errors="replace")
                sections = [{"heading": "Genel", "content": text, "index": 0}]
        except Exception as e:
            log_error(f"Section extraction hatası ({file_type}): {e}", "enhancer")
        
        return sections
    
    def _extract_pdf_sections(self, file_obj: BinaryIO) -> List[Dict[str, Any]]:
        """PDF bölümlerini çıkar"""
        try:
            from pypdf import PdfReader
        except ImportError:
            from PyPDF2 import PdfReader
        
        file_obj.seek(0)
        reader = PdfReader(file_obj)
        
        all_text = ""
        for page in reader.pages:
            text = page.extract_text() or ""
            all_text += text + "\n"
        
        return self._split_text_by_headings(all_text)
    
    def _extract_docx_sections(self, file_obj: BinaryIO) -> List[Dict[str, Any]]:
        """DOCX bölümlerini çıkar (Word stilleri ile) — paragraf aralığı da kaydedilir"""
        from docx import Document
        
        file_obj.seek(0)
        doc = Document(file_obj)
        
        sections = []
        current_heading = "Giriş"
        current_content = []
        idx = 0
        para_start = 0  # Bu bölümün başladığı paragraf index'i
        
        for para_idx, para in enumerate(doc.paragraphs):
            style_name = para.style.name if para.style else ""
            text = para.text.strip()
            
            if not text:
                current_content.append("")
                continue
            
            if style_name.startswith("Heading"):
                # Önceki bölümü kaydet
                if current_content:
                    sections.append({
                        "heading": current_heading,
                        "content": "\n".join(current_content),
                        "index": idx,
                        "para_start": para_start,
                        "para_end": para_idx - 1
                    })
                    idx += 1
                current_heading = text
                current_content = []
                para_start = para_idx + 1  # Heading'den sonraki paragraf başlangıcı
            else:
                current_content.append(text)
        
        # Son bölümü kaydet
        if current_content:
            sections.append({
                "heading": current_heading,
                "content": "\n".join(current_content),
                "index": idx,
                "para_start": para_start,
                "para_end": len(doc.paragraphs) - 1
            })
        
        if not sections:
            # Heading yoksa tüm içeriği tek bölüm olarak al
            full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            sections = [{
                "heading": "Genel", "content": full_text, "index": 0,
                "para_start": 0, "para_end": len(doc.paragraphs) - 1
            }]
        
        return sections
    
    def _extract_xlsx_sections(self, file_obj: BinaryIO) -> List[Dict[str, Any]]:
        """
        XLSX bölümlerini çıkar — v3.3.0: Veri bloğu bazlı bölümleme.
        Büyük sheet'ler boş satır gap'lerine göre alt-section'lara ayrılır.
        Header satırı her alt-section'a prefix olarak eklenir.
        """
        from openpyxl import load_workbook
        
        file_obj.seek(0)
        wb = load_workbook(file_obj, data_only=True)
        sections = []
        global_idx = 0
        global_row_counter = 0
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            
            # Tüm satırları oku
            all_rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(cell) if cell is not None else "" for cell in row]
                has_data = any(c.strip() for c in cells)
                all_rows.append((cells, has_data))
            
            if not all_rows:
                continue
            
            # Header satırını tespit et (ilk veri satırı, kısa metin, benzersiz değerler)
            header_text = ""
            header_row_idx = -1
            for ri, (cells, has_data) in enumerate(all_rows):
                if has_data:
                    non_empty = [c for c in cells if c.strip()]
                    all_short = all(len(c) < 50 for c in non_empty)
                    all_unique = len(set(non_empty)) == len(non_empty)
                    if non_empty and all_short and all_unique and len(non_empty) >= 2:
                        header_text = " | ".join(non_empty)
                        header_row_idx = ri
                    break
            
            # Veri satırlarını boş satır gap'lerine göre data bloklarına böl
            data_blocks = []
            current_block = []
            gap_count = 0
            
            start_row = header_row_idx + 1 if header_row_idx >= 0 else 0
            
            for ri in range(start_row, len(all_rows)):
                cells, has_data = all_rows[ri]
                if has_data:
                    if gap_count >= 2 and current_block:
                        # 2+ boş satır = yeni data bloğu
                        data_blocks.append(current_block)
                        current_block = []
                    current_block.append(" | ".join(cells))
                    gap_count = 0
                else:
                    gap_count += 1
            
            if current_block:
                data_blocks.append(current_block)
            
            # Her data bloğunu section olarak ekle
            if not data_blocks:
                continue
            
            # Tek blok varsa sheet adıyla section oluştur
            if len(data_blocks) == 1:
                block_text = "\n".join(data_blocks[0])
                if header_text:
                    block_text = f"[Başlıklar: {header_text}]\n{block_text}"
                
                para_start = global_row_counter
                para_end = global_row_counter + len(data_blocks[0]) - 1
                sections.append({
                    "heading": sheet_name,
                    "content": block_text,
                    "index": global_idx,
                    "para_start": para_start,
                    "para_end": para_end
                })
                global_idx += 1
                global_row_counter += len(data_blocks[0])
            else:
                # Birden fazla blok — her bloğu alt-section yap
                for bi, block in enumerate(data_blocks):
                    block_text = "\n".join(block)
                    if header_text:
                        block_text = f"[Başlıklar: {header_text}]\n{block_text}"
                    
                    block_heading = f"{sheet_name} — Bölüm {bi + 1}"
                    para_start = global_row_counter
                    para_end = global_row_counter + len(block) - 1
                    sections.append({
                        "heading": block_heading,
                        "content": block_text,
                        "index": global_idx,
                        "para_start": para_start,
                        "para_end": para_end
                    })
                    global_idx += 1
                    global_row_counter += len(block)
        
        return sections
    
    def _extract_csv_sections(self, file_obj: BinaryIO) -> List[Dict[str, Any]]:
        """
        CSV bölümlerini çıkar — v3.3.0: Satır bazlı bölümleme.
        CSVProcessor mantığıyla tutarlı: delimiter tespiti, header detection.
        """
        import csv as _csv
        
        file_obj.seek(0)
        raw = file_obj.read()
        
        # Encoding tespiti
        text = None
        try:
            from charset_normalizer import from_bytes
            result = from_bytes(raw).best()
            if result and result.encoding:
                text = str(result)
        except (ImportError, Exception):
            pass
        
        if text is None:
            for enc in ['utf-8', 'utf-8-sig', 'cp1254', 'iso-8859-9', 'latin-1']:
                try:
                    text = raw.decode(enc)
                    break
                except (UnicodeDecodeError, UnicodeError):
                    continue
            if text is None:
                text = raw.decode('utf-8', errors='replace')
        
        if not text or not text.strip():
            return []
        
        # Delimiter tespiti
        try:
            sample = "\n".join(text.split("\n")[:10])
            dialect = _csv.Sniffer().sniff(sample, delimiters=',;\t|')
            delimiter = dialect.delimiter
        except _csv.Error:
            counts = {',': text.count(','), ';': text.count(';'), '\t': text.count('\t')}
            delimiter = max(counts, key=counts.get)
        
        # CSV parse
        reader = _csv.reader(io.StringIO(text), delimiter=delimiter)
        all_rows = list(reader)
        
        if not all_rows:
            return []
        
        # Header tespiti
        has_header = False
        if len(all_rows) >= 2:
            first = all_rows[0]
            non_empty = [c for c in first if c.strip()]
            if non_empty and len(non_empty) >= 2:
                all_text = all(not c.strip().replace('.', '').replace(',', '').isdigit() for c in non_empty)
                all_short = all(len(c) < 50 for c in non_empty)
                all_unique = len(set(c.strip().lower() for c in non_empty)) == len(non_empty)
                has_header = all_text and all_short and all_unique
        
        header_row = all_rows[0] if has_header else None
        data_rows = all_rows[1:] if has_header else all_rows
        header_text = " | ".join(c.strip() for c in header_row if c.strip()) if header_row else ""
        
        # Bölümlere ayır (50 satır/bölüm)
        MAX_ROWS_PER_SECTION = 50
        sections = []
        
        for i in range(0, len(data_rows), MAX_ROWS_PER_SECTION):
            batch = data_rows[i:i + MAX_ROWS_PER_SECTION]
            row_texts = []
            for row in batch:
                row_text = " | ".join(c.strip() for c in row if c.strip())
                if row_text:
                    row_texts.append(row_text)
            
            if not row_texts:
                continue
            
            block_text = ""
            if header_text:
                block_text = f"[Başlıklar: {header_text}]\n"
            block_text += "\n".join(row_texts)
            
            section_heading = header_text or f"CSV Veri Bloğu {len(sections) + 1}"
            if len(sections) > 0 or i > 0:
                section_heading = f"{section_heading} — Bölüm {len(sections) + 1}"
            
            sections.append({
                "heading": section_heading,
                "content": block_text,
                "index": len(sections),
                "para_start": i,
                "para_end": i + len(row_texts) - 1
            })
        
        return sections
    
    def _extract_pptx_sections(self, file_obj: BinaryIO) -> List[Dict[str, Any]]:
        """PPTX her slaytı bir bölüm olarak çıkar — paragraf aralığı da kaydedilir"""
        from pptx import Presentation
        
        file_obj.seek(0)
        prs = Presentation(file_obj)
        sections = []
        global_para_counter = 0
        
        for idx, slide in enumerate(prs.slides):
            texts = []
            slide_title = f"Slayt {idx + 1}"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if shape.shape_type is not None and shape.shape_type == 13:
                        continue  # Picture
                    if shape == slide.shapes.title:
                        slide_title = shape.text.strip()
                    else:
                        texts.append(shape.text.strip())
            
            if texts:
                para_start = global_para_counter
                para_end = global_para_counter + len(texts) - 1
                sections.append({
                    "heading": slide_title,
                    "content": "\n".join(texts),
                    "index": idx,
                    "para_start": para_start,
                    "para_end": para_end
                })
                global_para_counter += len(texts)
        
        return sections
    
    def _extract_txt_sections(self, file_obj: BinaryIO) -> List[Dict[str, Any]]:
        """TXT dosyasını satır bazlı bölümlere ayır"""
        file_obj.seek(0)
        text = file_obj.read().decode("utf-8", errors="replace")
        return self._split_text_by_headings(text)
    
    def _split_text_by_headings(self, text: str) -> List[Dict[str, Any]]:
        """Metni heading pattern'lerine göre bölümlere ayır — paragraf aralığı da kaydedilir"""
        lines = text.split("\n")
        sections = []
        current_heading = "Giriş"
        current_content = []
        idx = 0
        line_start = 0  # Bu bölümün başladığı satır index'i
        
        heading_patterns = [
            r'^\d+[\.\)]\s+\S',
            r'^[A-ZÇĞİÖŞÜ][A-ZÇĞİÖŞÜ\s]{3,}$',
            r'^(?:BÖLÜM|MADDE|KISIM|BAŞLIK)\s',
        ]
        
        for line_idx, line in enumerate(lines):
            stripped = line.strip()
            if not stripped:
                current_content.append("")
                continue
            
            is_heading = False
            if len(stripped) < 80:
                for pattern in heading_patterns:
                    if re.match(pattern, stripped):
                        is_heading = True
                        break
            
            if is_heading:
                if current_content:
                    sections.append({
                        "heading": current_heading,
                        "content": "\n".join(current_content),
                        "index": idx,
                        "para_start": line_start,
                        "para_end": line_idx - 1
                    })
                    idx += 1
                current_heading = stripped
                current_content = []
                line_start = line_idx + 1  # Heading'den sonraki satır
            else:
                current_content.append(stripped)
        
        if current_content:
            sections.append({
                "heading": current_heading,
                "content": "\n".join(current_content),
                "index": idx,
                "para_start": line_start,
                "para_end": len(lines) - 1
            })
        
        # Hiç bölüm bulunmadıysa tüm metni tek bölüm yap
        if not sections:
            sections = [{
                "heading": "Genel", "content": text.strip(), "index": 0,
                "para_start": 0, "para_end": len(lines) - 1
            }]
        
        return sections
    
    # ─────────────────────────────────────────
    #  ADIM 2: CatBoost Priority Analizi
    # ─────────────────────────────────────────
    
    def _catboost_prioritize(
        self,
        sections: List[Dict[str, Any]],
        file_type: str,
        violations: List[Dict[str, Any]]
    ) -> Dict[str, Any]:
        """
        CatBoost ile her bölümün RAG kalitesini tahmin et.
        Model yoksa maturity violation'lara göre basit scoring.
        """
        
        catboost = self._get_catboost()
        fe = self._get_feature_extractor()
        
        section_priorities = []
        has_catboost = catboost is not None and catboost.is_ready() and fe is not None
        
        # Violation → section eşleştirme
        violation_names = [v.get("name", "") for v in violations]
        
        # v3.2.1: Debug — hangi violation'lar tespit edildi
        log_system_event("DEBUG", 
            f"CatBoost prioritize: {len(sections)} section, "
            f"violations={violation_names}, catboost={has_catboost}", 
            "enhancer")
        
        for section in sections:
            content = section.get("content", "")
            heading = section.get("heading", "")
            word_count = len(content.split())
            
            if has_catboost:
                # CatBoost ile kalite tahmini
                try:

                    
                    # Fake RAG result oluştur (feature extractor'ın beklediği format)
                    fake_result = {
                        "chunk_id": 0,
                        "content": content,
                        "chunk_text": content,
                        "score": 0.5,
                        "quality_score": 0.5,
                        "metadata": {
                            "heading": heading,
                            "file_type": f".{file_type.lower()}"
                        },
                        "file_type": f".{file_type.lower()}"
                    }
                    
                    feature_matrix, _ = fe.build_feature_matrix(
                        results=[fake_result],
                        user_id=None,
                        query=heading or "doküman kalite analizi"
                    )
                    
                    if feature_matrix is not None and len(feature_matrix) > 0:
                        scores = catboost.predict(feature_matrix)
                        priority = 1.0 - float(scores[0])  # Düşük skor = yüksek iyileştirme önceliği
                    else:
                        priority = self._heuristic_priority(content, heading, violation_names)
                except Exception as e:
                    log_warning(f"CatBoost prediction hatası: {e}", "enhancer")
                    priority = self._heuristic_priority(content, heading, violation_names)
            else:
                # CatBoost yoksa heuristic
                priority = self._heuristic_priority(content, heading, violation_names)
            
            weakness_types = self._detect_weaknesses(content, heading, violation_names)
            
            section_priorities.append({
                "index": section["index"],
                "heading": heading,
                "priority": round(priority, 3),
                "weakness_types": weakness_types,
                "word_count": word_count,
                "catboost_used": has_catboost
            })
        
        # Özet
        avg_priority = sum(s["priority"] for s in section_priorities) / max(len(section_priorities), 1)
        high_priority_count = sum(1 for s in section_priorities if s["priority"] > 0.5)
        
        return {
            "sections": section_priorities,
            "summary": {
                "total_sections": len(section_priorities),
                "high_priority_count": high_priority_count,
                "average_priority": round(avg_priority, 3),
                "catboost_available": has_catboost,
                "violation_count": len(violations)
            }
        }
    
    def _heuristic_priority(
        self,
        content: str,
        heading: str,
        violation_names: List[str]
    ) -> float:
        """
        CatBoost yoksa heuristic ile priority hesapla — tüm dosya türleri.
        v3.3.0 [C3]: Min-max normalizasyon — CatBoost ile tutarlı 0-1 dağılım.
        """
        raw_score = 0.0  # Ham skor (normalize edilecek)
        
        word_count = len(content.split())
        
        # Çok kısa içerik → yüksek priority
        if word_count < 20:
            raw_score += 3.0
        elif word_count < 50:
            raw_score += 1.5
        
        # Heading yoksa veya Generic ise
        if not heading or heading in ("Genel", "Giriş"):
            raw_score += 2.0
        
        # Violation eşleştirme — PDF/DOCX/TXT
        if "Başlık Hiyerarşisi" in violation_names or "Word Stilleri" in violation_names:
            raw_score += 1.5
        if "Metin Yoğunluğu" in violation_names or "Metin İçeriği" in violation_names:
            raw_score += 1.0
        if "Tablo Formatı" in violation_names:
            raw_score += 1.0
        if "Türkçe Karakter" in violation_names:
            raw_score += 1.0
        if "Gereksiz İçerik" in violation_names or "Gereksiz Boşluklar" in violation_names:
            raw_score += 1.0
        
        # v3.2.1: Excel'e özel ihlal boost
        if "İlk Satır Başlık" in violation_names:
            raw_score += 1.5
        if "Merge Hücreler" in violation_names:
            raw_score += 1.0
        if "Açıklama Satırları" in violation_names:
            raw_score += 1.0
        if "Boş Satır/Sütun" in violation_names:
            raw_score += 1.0
        if "Formül vs Değer" in violation_names:
            raw_score += 1.0
        if "Tutarlı Veri Tipi" in violation_names:
            raw_score += 0.5
        if "Gizli Sheet" in violation_names:
            raw_score += 0.5
        
        # DOCX özel
        if "Metin Kutusu" in violation_names:
            raw_score += 1.5
        if "Liste Formatı" in violation_names:
            raw_score += 0.5
        
        # PPTX özel (v3.3.0: yeni kurallar)
        if "Slayt Başlıkları" in violation_names:
            raw_score += 1.0
        if "Speaker Notes" in violation_names:
            raw_score += 0.8
        if "Görsel/Metin Oranı" in violation_names:
            raw_score += 0.8
        if "Slayt Sayısı" in violation_names:
            raw_score += 0.5
        
        # v3.3.0 [C3]: Min-max normalizasyon
        # Olası max ham skor ≈ 15 (tüm violation'lar + kısa metin + heading yok)
        MAX_RAW_SCORE = 15.0
        normalized = min(raw_score / MAX_RAW_SCORE, 1.0)
        
        # Minimum 0.1 floor (tamamen sorunsuz bile olsa bir baseline sağlar)
        return max(round(normalized, 3), 0.1)
    
    def _detect_weaknesses(
        self,
        content: str,
        heading: str,
        violation_names: List[str]
    ) -> List[str]:
        """Bölüm bazında zayıflıkları tespit et — tüm dosya türleri desteklenir."""
        weaknesses = []
        
        word_count = len(content.split())
        
        if not heading or heading in ("Genel", "Giriş"):
            weaknesses.append("heading_missing")
        
        if word_count < 20:
            weaknesses.append("content_too_short")
        
        # Tablo kontrol
        if "|" in content and content.count("|") > 5:
            if "Tablo Formatı" in violation_names:
                weaknesses.append("table_format_issue")
        
        # Türkçe karakter
        if "Türkçe Karakter" in violation_names:
            bad_chars = ['Ã¼', 'Ã§', 'Ã¶', 'Ä±', 'ÅŸ']
            if any(bc in content for bc in bad_chars):
                weaknesses.append("encoding_issue")
        
        # PDF / DOCX / TXT yapısal sorunlar
        if "Başlık Hiyerarşisi" in violation_names:
            weaknesses.append("structure_weak")
        if "Metin Yoğunluğu" in violation_names:
            weaknesses.append("low_density")
        if "Gereksiz İçerik" in violation_names:
            weaknesses.append("redundant_content")
        if "Gereksiz Boşluklar" in violation_names:
            weaknesses.append("excess_whitespace")
        
        # DOCX özel
        if "Word Stilleri" in violation_names:
            weaknesses.append("structure_weak")
        if "Metin Kutusu" in violation_names:
            weaknesses.append("textbox_issue")
        if "Liste Formatı" in violation_names:
            weaknesses.append("list_format_issue")
        
        # v3.2.1: Excel'e özel ihlal eşleştirmesi
        if "İlk Satır Başlık" in violation_names:
            weaknesses.append("header_row_missing")
        if "Merge Hücreler" in violation_names:
            weaknesses.append("merged_cells")
        if "Açıklama Satırları" in violation_names:
            weaknesses.append("description_rows")
        if "Boş Satır/Sütun" in violation_names:
            weaknesses.append("empty_gaps")
        if "Tutarlı Veri Tipi" in violation_names:
            weaknesses.append("inconsistent_types")
        if "Formül vs Değer" in violation_names:
            weaknesses.append("formula_issue")
        if "Gizli Sheet" in violation_names:
            weaknesses.append("hidden_sheets")
        
        # PPTX özel
        if "Metin İçeriği" in violation_names:
            weaknesses.append("low_density")
        if "Slayt Başlıkları" in violation_names:
            weaknesses.append("heading_missing")
        
        return weaknesses
    
    # ─────────────────────────────────────────
    #  ADIM 3: LLM İyileştirme
    # ─────────────────────────────────────────
    
    def _llm_enhance(
        self,
        sections: List[Dict[str, Any]],
        violations: List[Dict[str, Any]],
        file_type: str,
        catboost_analysis: Dict[str, Any],
        progress_callback=None
    ) -> List[EnhancedSection]:
        """LLM ile düşük kaliteli bölümleri iyileştir"""
        
        enhanced_sections = []
        cb_sections = catboost_analysis.get("sections", [])
        
        for section in sections:
            idx = section["index"]
            heading = section.get("heading", "")
            content = section.get("content", "")
            
            # CatBoost priority bul
            cb_info = next((s for s in cb_sections if s["index"] == idx), None)
            priority = cb_info["priority"] if cb_info else 0.5
            weakness_types = cb_info.get("weakness_types", []) if cb_info else []
            
            # v3.2.1: Debug loglama — karar sürecini izlemek için
            log_system_event("DEBUG", 
                f"Bölüm [{idx}] '{heading[:30]}': priority={priority:.3f}, "
                f"weaknesses={weakness_types}, catboost={cb_info is not None}", 
                "enhancer")
            
            # Priority düşükse (yani kalite yüksek) VE zayıflık yoksa değiştirme
            # v3.2.1: weakness_types varsa priority ne olursa olsun LLM'e gönder
            if priority < 0.4 and not weakness_types:
                log_system_event("DEBUG", 
                    f"Bölüm [{idx}] → SKIP (priority={priority:.3f}, no weaknesses)", 
                    "enhancer")
                enhanced_sections.append(EnhancedSection(
                    section_index=idx,
                    heading=heading,
                    original_text=content,
                    enhanced_text=content,
                    change_type="no_change",
                    explanation="Bu bölüm yeterli kalitede.",
                    priority=priority
                ))
                # v3.3.0 [C5]: Progress callback — "Bölüm X/Y atlandı"
                if progress_callback:
                    try:
                        progress_callback(idx + 1, len(sections), heading, "skipped")
                    except Exception:
                        pass
                continue
            
            # v3.3.0 [C5]: Progress callback — "Bölüm X/Y iyileştiriliyor..."
            if progress_callback:
                try:
                    progress_callback(idx + 1, len(sections), heading, "processing")
                except Exception:
                    pass
            
            # LLM ile iyileştir — v3.3.3 [Faz 2]: max 2 retry desteği
            MAX_ENHANCEMENT_RETRIES = 2
            try:
                llm_result = self._call_llm_for_enhancement(
                    heading=heading,
                    content=content,
                    weakness_types=weakness_types,
                    violations=violations,
                    file_type=file_type
                )
                
                enhanced_text = llm_result.get("enhanced_text", content)
                
                # ─── Content Integrity Validation ───
                validator = get_integrity_validator()
                integrity = validator.validate(
                    original=content,
                    enhanced=enhanced_text,
                    file_type=file_type,
                    weakness_types=weakness_types
                )
                
                # 🛡️ v3.3.3 [Faz 2]: Retry with Corrective Prompt
                # Bütünlük başarısızsa → düzeltici prompt ile tekrar dene
                retry_attempt = 0
                while not integrity.is_valid and retry_attempt < MAX_ENHANCEMENT_RETRIES:
                    retry_attempt += 1
                    log_system_event(
                        "INFO",
                        f"Bölüm [{idx}] integrity fail (skor={integrity.score:.3f}), "
                        f"retry {retry_attempt}/{MAX_ENHANCEMENT_RETRIES}...",
                        "enhancer"
                    )
                    
                    # Progress callback: retrying durumu
                    if progress_callback:
                        try:
                            progress_callback(
                                idx + 1, len(sections), heading,
                                f"retrying_{retry_attempt}"
                            )
                        except Exception:
                            pass
                    
                    # Düzeltici prompt ile LLM'i tekrar çağır
                    corrective_result = self._call_llm_corrective(
                        heading=heading,
                        content=content,
                        failed_enhanced_text=enhanced_text,
                        integrity_issues=integrity.issues,
                        lost_entities=integrity.lost_entities,
                        hallucinated_entities=integrity.hallucinated_entities,
                        weakness_types=weakness_types,
                        violations=violations,
                        file_type=file_type
                    )
                    enhanced_text = corrective_result.get("enhanced_text", content)
                    llm_result = corrective_result  # heading vb. güncellemesi için
                    
                    # Yeniden validate et
                    integrity = validator.validate(
                        original=content,
                        enhanced=enhanced_text,
                        file_type=file_type,
                        weakness_types=weakness_types
                    )
                
                if not integrity.is_valid:
                    # Tüm denemeler başarısız → Orijinali koru
                    total_attempts = retry_attempt + 1
                    log_warning(
                        f"Bölüm [{idx}] {total_attempts} denemede de bütünlük BAŞARISIZ "
                        f"(skor={integrity.score:.3f}): {'; '.join(integrity.issues[:2])}",
                        "enhancer"
                    )
                    enhanced_sections.append(EnhancedSection(
                        section_index=idx,
                        heading=heading,  # Orijinal başlık koru
                        original_text=content,
                        enhanced_text=content,  # ORİJİNALİ KORU
                        change_type="integrity_failed",
                        explanation=(
                            f"İyileştirme reddedildi ({total_attempts} deneme) — "
                            f"{'; '.join(integrity.issues[:2])}. "
                            f"Bütünlük skoru: {integrity.score:.0%}"
                        ),
                        priority=priority,
                        violations=weakness_types,
                        integrity_score=integrity.score,
                        integrity_issues=integrity.issues + integrity.warnings
                    ))
                    if progress_callback:
                        try:
                            progress_callback(idx + 1, len(sections), heading, "error")
                        except Exception:
                            pass
                    continue
                
                # Bütünlük OK — iyileştirmeyi kabul et
                if retry_attempt > 0:
                    log_system_event(
                        "INFO",
                        f"Bölüm [{idx}] retry {retry_attempt} sonrası integrity BAŞARILI "
                        f"(skor={integrity.score:.3f})",
                        "enhancer"
                    )
                
                enhanced_sections.append(EnhancedSection(
                    section_index=idx,
                    heading=llm_result.get("heading", heading),
                    original_text=content,
                    enhanced_text=enhanced_text,
                    change_type=llm_result.get("change_type", "content_restructured"),
                    explanation=(
                        llm_result.get("explanation", "LLM ile iyileştirildi.")
                        + (f" (retry {retry_attempt} sonrası başarılı)" if retry_attempt > 0 else "")
                    ),
                    priority=priority,
                    violations=weakness_types,
                    integrity_score=integrity.score,
                    integrity_issues=integrity.warnings  # Sadece uyarılar
                ))
                if progress_callback:
                    try:
                        progress_callback(idx + 1, len(sections), heading, "processing")
                    except Exception:
                        pass
                
            except Exception as e:
                log_warning(f"LLM enhancement hatası (bölüm {idx}): {e}", "enhancer")
                # v3.2.1: LLM hatası kullanıcıya gösterilmeli — no_change yerine llm_error
                error_msg = str(e)[:150]
                if "Failed to resolve" in error_msg or "getaddrinfo" in error_msg:
                    user_msg = "LLM API sunucusuna bağlanılamadı (DNS hatası). Ağ bağlantınızı kontrol edin."
                elif "timed out" in error_msg.lower() or "timeout" in error_msg.lower():
                    user_msg = "LLM API yanıt zaman aşımı. Lütfen tekrar deneyin."
                elif "401" in error_msg or "403" in error_msg:
                    user_msg = "LLM API yetkilendirme hatası. API anahtarını kontrol edin."
                else:
                    user_msg = f"LLM iyileştirme yapılamadı: {error_msg}"
                
                enhanced_sections.append(EnhancedSection(
                    section_index=idx,
                    heading=heading,
                    original_text=content,
                    enhanced_text=content,
                    change_type="llm_error",
                    explanation=user_msg,
                    priority=priority,
                    violations=weakness_types
                ))
                # v3.3.0 [C5]: Progress callback — hata sonrası da bildirim
                if progress_callback:
                    try:
                        progress_callback(idx + 1, len(sections), heading, "error")
                    except Exception:
                        pass
        
        return enhanced_sections
    
    def _call_llm_for_enhancement(
        self,
        heading: str,
        content: str,
        weakness_types: List[str],
        violations: List[Dict[str, Any]],
        file_type: str
    ) -> Dict[str, Any]:
        """Tek bir bölüm için LLM'e iyileştirme isteği gönder"""
        from app.core.llm import call_llm_api
        
        # v3.3.0: Akıllı kırpma — heading/ayıraç bazlı kesim noktası
        max_content_chars = 6000
        remaining_text = ""
        context_overlap = ""  # LLM'e bağlam olarak gönderilecek kırpılan kısmın son 300 karakteri
        
        if len(content) > max_content_chars:
            # Öncelik 1: Heading veya bölüm ayracından kes
            heading_cut_patterns = [
                r'\n\d+[\.\)]\s+\S',       # 1. veya 1) ile başlayan satır
                r'\n[A-ZÇĞİÖŞÜ]{4,}',     # TAMAMEN BÜYÜK HARF satır
                r'\n---+',                   # --- ayıraç
                r'\n#{1,3}\s',              # Markdown heading
                r'\n(?:BÖLÜM|MADDE|KISIM)\s', # Türkçe bölüm kelimeleri
            ]
            
            best_cut = -1
            # max_content_chars aralığında heading ara (son %30'luk dilimde)
            search_start = int(max_content_chars * 0.7)
            search_end = max_content_chars
            
            for pattern in heading_cut_patterns:
                matches = list(re.finditer(pattern, content[search_start:search_end]))
                if matches:
                    # En son eşleşmeyi al (mümkün olduğunca fazla içerik gönder)
                    best_cut = search_start + matches[-1].start()
                    break
            
            # Öncelik 2: Paragraf sınırından kes
            if best_cut < 0:
                cut_point = content.rfind("\n\n", 0, max_content_chars)
                if cut_point > max_content_chars // 2:
                    best_cut = cut_point
            
            # Öncelik 3: Tek satır sonundan kes
            if best_cut < 0:
                cut_point = content.rfind("\n", 0, max_content_chars)
                if cut_point > max_content_chars // 2:
                    best_cut = cut_point
            
            # Son çare: karakter limitinden kes
            if best_cut < 0:
                best_cut = max_content_chars
            
            truncated = content[:best_cut]
            remaining_text = content[best_cut:]
            
            # v3.3.0: Kalan kısmın ilk 300 karakterini LLM'e context olarak gönder
            # Bu sayede LLM paragrafı yarıda bırakmaz
            context_overlap = remaining_text[:300].strip()
            
            log_system_event(
                "INFO",
                f"İçerik kırpıldı: {len(content)} → {len(truncated)} karakter "
                f"(heading/ayıraç bazlı kesim, kalan {len(remaining_text)} karakter orijinal korunacak)",
                "enhancer"
            )
        else:
            truncated = content
        
        # 🛡️ v3.3.3: Anchor Extraction — kritik verileri placeholder ile koru
        from app.services.content_anchor_service import get_anchor_service
        anchor_svc = get_anchor_service()
        anchor_result = anchor_svc.extract_anchors(truncated)
        anchored_text = anchor_result.sanitized_text
        anchor_registry = anchor_result.anchor_registry
        
        if anchor_result.anchor_count > 0:
            log_system_event(
                "INFO",
                f"Anchor extraction: {anchor_result.anchor_count} kritik veri korundu — {anchor_result.anchor_types}",
                "enhancer"
            )
        
        # Weakness'e göre talimat oluştur
        instructions = self._build_fix_instructions(weakness_types)
        
        violation_list = "\n".join(
            f"- {v.get('name', '')}: {v.get('detail', '')}" 
            for v in violations if v.get("status") in ("fail", "warning")
        ) or "Belirgin ihlal yok."
        
        # v3.3.0: Context overlap bilgisini prompt'a ekle
        context_note = ""
        if context_overlap:
            context_note = f"\n\n**NOT:** Bu bölümün devamı aşağıdaki gibidir (sadece bağlam için, iyileştirmeye DAHİL ETMEYİN):\n{context_overlap}...\n"
        
        # 🛡️ v3.3.3 [Faz 3]: Structured Prompt — Fenced Critical Blocks
        anchor_note = ""
        if anchor_registry:
            # Anchor ID'lerini listele — LLM hangi placeholder'ların dokunulmaz olduğunu bilsin
            anchor_ids_display = "  ".join(anchor_registry.keys())
            anchor_type_summary = ", ".join(
                f"{atype}: {count}" for atype, count in anchor_result.anchor_types.items()
            )
            anchor_note = (
                f"\n\n**⛔ DOKUNMA BÖLGESİ — FROZEN DATA ({anchor_result.anchor_count} adet, {anchor_type_summary}):**\n"
                "---FROZEN_START---\n"
                f"{anchor_ids_display}\n"
                "---FROZEN_END---\n\n"
                "**KURALLAR:**\n"
                "- Yukarıdaki ‹‹ANC_XXX›› placeholder'lar gerçek verileri (sayı, tarih, URL vb.) temsil eder.\n"
                "- Bu placeholder'ları AYNEN KORU — SİLME, DEĞİŞTİRME, YERİNİ DEĞİŞTİRME.\n"
                "- Placeholder'ın etrafındaki bağlam cümlesini değiştirirsen bile placeholder'ı KORU.\n"
            )
        
        prompt = f"""Sen bir doküman optimize uzmanısın. Aşağıdaki doküman bölümünü RAG (Retrieval-Augmented Generation) sistemi için optimize et.

**Dosya Tipi:** {file_type}
**Bölüm Başlığı:** {heading or '(Başlık yok)'}
**Tespit Edilen Sorunlar:** {', '.join(weakness_types) if weakness_types else 'Genel kalite iyileştirmesi'}

**İyileştirme Talimatları:**
{instructions}

**Maturity İhlalleri:**
{violation_list}{anchor_note}

**Orijinal İçerik:**
{anchored_text}{context_note}

**Yanıt Formatı (JSON):**
{{
    "heading": "İyileştirilmiş başlık (yoksa uygun bir başlık öner)",
    "enhanced_text": "İyileştirilmiş metin (orijinal yapıyı koru, sadece sorunları düzelt)",
    "change_type": "heading_added|content_restructured|table_fixed|encoding_fixed|formatting_improved",
    "explanation": "Yapılan değişikliklerin kısa açıklaması (Türkçe)"
}}

**KRİTİK KURALLAR:**
- Orijinal metnin ANLAMINI ve İÇERİĞİNİ koru
- Yapısal iyileştirme yap: başlık ekleme, paragraf düzenleme, format düzeltme, bölümlendirme
- İçerik iyileştirme yap: yazım/imla düzeltme, cümle netliği artırma, tutarsızlıkları giderme
- ASLA kendi yorumunu, açıklamanı veya ek bilgi EKLEME
- ASLA "Not:", "Açıklama:", "Yorum:", "Dipnot:" gibi meta-tekstler ekleme
- ASLA orijinal bilgiyi silme — sadece daha iyi ifade et
- Yanıtı SADECE JSON formatında ver, başka hiçbir metin ekleme"""

        messages = [
            {"role": "system", "content": "Sen doküman optimizasyon uzmanısın. Yanıtlarını SADECE JSON formatında ver."},
            {"role": "user", "content": prompt}
        ]
        
        response = call_llm_api(messages)
        
        # JSON parse
        result = self._parse_llm_response(response, heading, content)
        
        # 🛡️ v3.3.3: Anchor Re-injection + Recovery
        if anchor_registry:
            enhanced_text = result.get("enhanced_text", "")
            # Adım 1: Placeholder'ları orijinal değerlerle değiştir
            enhanced_text = anchor_svc.reinject_anchors(enhanced_text, anchor_registry)
            # Adım 2: LLM'in sildiği anchor'ları kurtarıp geri ekle
            enhanced_text, recovered = anchor_svc.recover_missing(
                enhanced_text, anchor_registry, truncated
            )
            result["enhanced_text"] = enhanced_text
            
            if recovered:
                result["explanation"] = (
                    result.get("explanation", "") +
                    f" (⚠ {len(recovered)} kayıp veri otomatik kurtarıldı)"
                )
        
        # v3.3.0: Kırpılan içeriği enhanced text'e ekle — yumuşak geçiş
        if remaining_text:
            enhanced_part = result.get("enhanced_text", "")
            # Birleşim noktasında çift newline ile temiz geçiş sağla
            result["enhanced_text"] = enhanced_part.rstrip() + "\n\n" + remaining_text.lstrip()
            result["explanation"] = (result.get("explanation", "") + 
                f" (İlk {len(truncated)} karakter iyileştirildi, "
                f"kalan {len(remaining_text)} karakter orijinal olarak korundu)")
        
        return result
    
    # 🛡️ v3.3.3 [Faz 2]: Düzeltici Prompt ile LLM Retry
    def _call_llm_corrective(
        self,
        heading: str,
        content: str,
        failed_enhanced_text: str,
        integrity_issues: List[str],
        lost_entities: List[str],
        hallucinated_entities: List[str],
        weakness_types: List[str],
        violations: List[Dict[str, Any]],
        file_type: str
    ) -> Dict[str, Any]:
        """
        Düzeltici prompt ile LLM'i tekrar çağırır.
        
        Önceki iyileştirme integrity doğrulamasını geçemediyse,
        hataları ve kayıp verileri LLM'e göstererek düzeltme ister.
        Anchor extraction + re-injection dahildir.
        
        Args:
            heading: Bölüm başlığı
            content: Orijinal metin (kaynak)
            failed_enhanced_text: Başarısız iyileştirilmiş metin
            integrity_issues: Tespit edilen bütünlük sorunları
            lost_entities: Kaybolan varlıklar listesi
            hallucinated_entities: Halüsinasyon şüphesi olan varlıklar
            weakness_types: Bilinen zayıflık tipleri
            violations: Maturity ihlalleri
            file_type: Dosya tipi
            
        Returns:
            Dict: LLM sonucu ({heading, enhanced_text, change_type, explanation})
        """
        from app.core.llm import call_llm_api
        from app.services.content_anchor_service import get_anchor_service
        
        # Anchor extraction — orijinal metinden (kırpılmış)
        # v3.3.3 [Code Review Fix]: Uzun içerikler için kırpma uygula
        max_content_chars = 6000
        if len(content) > max_content_chars:
            # Basit kırpma — retry'da karmaşık heading kesimi gereksiz
            cut_point = content.rfind("\n\n", 0, max_content_chars)
            if cut_point < max_content_chars // 2:
                cut_point = content.rfind("\n", 0, max_content_chars)
            if cut_point < max_content_chars // 2:
                cut_point = max_content_chars
            working_content = content[:cut_point]
        else:
            working_content = content
        
        anchor_svc = get_anchor_service()
        anchor_result = anchor_svc.extract_anchors(working_content)
        anchored_text = anchor_result.sanitized_text
        anchor_registry = anchor_result.anchor_registry
        
        # Hata özetini oluştur
        issues_text = "\n".join(f"- {issue}" for issue in integrity_issues[:5])
        
        lost_text = ", ".join(lost_entities[:10]) if lost_entities else "Belirtilmedi"
        
        hallucinated_text = ""
        if hallucinated_entities:
            hallucinated_text = (
                f"\n\n**EKLENEN SAHTE VERİLER (BUNLARI KALDIRMALISIN):**\n"
                f"{', '.join(hallucinated_entities[:5])}"
            )
        
        # Anchor bilgisi
        anchor_note = ""
        if anchor_registry:
            anchor_note = (
                "\n\n**⛔ KRİTİK VERİ KORUMA:**\n"
                "Metindeki ‹‹ANC_XXX›› formatındaki placeholder'lar gerçek verileri temsil eder.\n"
                "Bu placeholder'ları AYNEN KORU, SİLME, DEĞİŞTİRME. Yerlerini değiştirme.\n"
            )
        
        prompt = f"""⚠️ ÖNCEKİ İYİLEŞTİRMENDE CİDDİ HATALAR TESPİT EDİLDİ.
Lütfen orijinal metni BAZ ALARAK tekrar iyileştir. Bu sefer aşağıdaki hataları YAPMA.

**TESPİT EDİLEN HATALAR:**
{issues_text}

**KAYIP/DEĞİŞTİRİLEN VERİLER:**
{lost_text}{hallucinated_text}

**ORİJİNAL İÇERİK (BUNU BAZ AL):**
{anchored_text}{anchor_note}

**Dosya Tipi:** {file_type}
**Bölüm Başlığı:** {heading or '(Başlık yok)'}

**KRİTİK KURALLAR (BU SEFER MUTLAKA UYULMALIDIR):**
1. Orijinal metindeki TÜM sayıları (1.234,56 gibi) AYNEN koru — format değiştirme
2. TÜM tarihleri, URL'leri, email adreslerini AYNEN koru
3. ASLA orijinalde olmayan yeni bilgi, sayı veya isim EKLEME
4. ASLA "Not:", "Açıklama:", "Yorum:" gibi meta-tekstler ekleme
5. Sadece yapısal iyileştirme yap: başlık, paragraf, format düzenleme
6. ‹‹ANC_XXX›› placeholder'larını AYNEN BIRAK, silme

**Yanıt Formatı (JSON):**
{{
    "heading": "İyileştirilmiş başlık",
    "enhanced_text": "İyileştirilmiş metin",
    "change_type": "content_restructured",
    "explanation": "Yapılan düzeltmelerin açıklaması"
}}"""
        
        messages = [
            {
                "role": "system",
                "content": (
                    "Sen doküman optimizasyon uzmanısın. ÖNCEKİ denemende bütünlük hataları yapıldı. "
                    "Bu sefer ÇOK DİKKATLİ ol — orijinal verileri ASLA silme veya değiştirme. "
                    "Yanıtını SADECE JSON formatında ver."
                )
            },
            {"role": "user", "content": prompt}
        ]
        
        response = call_llm_api(messages)
        result = self._parse_llm_response(response, heading, content)
        
        # Anchor re-injection + recovery
        if anchor_registry:
            enhanced_text = result.get("enhanced_text", "")
            # Adım 1: Placeholder → orijinal değer
            enhanced_text = anchor_svc.reinject_anchors(enhanced_text, anchor_registry)
            # Adım 2: Kayıp anchor'ları kurtarıp geri ekle
            enhanced_text, recovered = anchor_svc.recover_missing(
                enhanced_text, anchor_registry, working_content
            )
            result["enhanced_text"] = enhanced_text
            
            if recovered:
                result["explanation"] = (
                    result.get("explanation", "") +
                    f" (⚠ retry: {len(recovered)} kayıp veri otomatik kurtarıldı)"
                )
        
        log_system_event(
            "DEBUG",
            f"Corrective LLM çağrısı tamamlandı: heading='{heading[:30]}'",
            "enhancer"
        )
        
        return result
    
    def _build_fix_instructions(self, weakness_types: List[str]) -> str:
        """Weakness tipine göre spesifik iyileştirme talimatları"""
        instructions = []
        
        instruction_map = {
            "heading_missing": "- Bu bölümün uygun bir başlığı yok. İçeriğe uygun, kısa ve açıklayıcı bir Heading ekle.",
            "content_too_short": "- İçerik çok kısa. Mevcut bilgiyi daha açıklayıcı hale getir (ama yeni bilgi EKLEME).",
            "table_format_issue": "- Tablo formatı bozuk. Sütun başlıklarını netleştir, hücre hizalamasını düzelt.",
            "encoding_issue": "- Türkçe karakter encoding sorunu var. Bozuk karakterleri (Ã¼→ü, Ã§→ç, Ã¶→ö, Ä±→ı, ÅŸ→ş) düzelt.",
            "structure_weak": "- Başlık hiyerarşisi zayıf. Alt başlıklar (##, ###) ekleyerek içeriği daha iyi yapılandır.",
            "low_density": "- Metin yoğunluğu düşük. Gereksiz boşlukları kaldır ve paragrafları birleştir.",
            "redundant_content": "- Tekrarlayan header/footer içerik var. Her bölümde tekrar eden boilerplate metinleri kaldır.",
            "excess_whitespace": "- Fazla boş satır/paragraf var. Gereksiz boşlukları temizle.",
            # Excel özel
            "header_row_missing": "- Excel'de başlık satırı eksik. İlk satırı sütun başlıkları olarak yeniden düzenle.",
            "merged_cells": "- Birleştirilmiş hücreler var. Her hücrenin kendi değerini taşıyacak şekilde yapıyı düzelt.",
            "description_rows": "- Veri tablosunun üstünde açıklama satırları var. Açıklamaları ayrı bir bölüme taşı veya kaldır.",
            "empty_gaps": "- Veri blokları arasında boş satırlar var. Boşlukları kaldırarak sürekli bir veri akışı oluştur.",
            "inconsistent_types": "- Sütunlarda karışık veri tipleri var. Her sütunda tutarlı format kullanılmalı.",
            "formula_issue": "- Formüllü hücreler var. Formülleri değerlerine çevir (hesaplanmış sonuçları yaz).",
            "hidden_sheets": "- Gizli sayfalar var. Gizli sayfa içeriğini ana bölüme dahil et veya gereksizse kaldır.",
            # DOCX özel
            "textbox_issue": "- Metin kutuları (text box) var. İçindeki metinleri normal paragraflara dönüştür.",
            "list_format_issue": "- Manuel liste kullanılmış. Word'ün yerleşik liste stillerini kullan.",
            # PPTX özel
        }
        
        for wt in weakness_types:
            if wt in instruction_map:
                instructions.append(instruction_map[wt])
        
        if not instructions:
            instructions.append("- Genel kalite iyileştirmesi: başlık yapısı, paragraf düzeni ve okunabilirliği artır.")
        
        return "\n".join(instructions)
    
    def _parse_llm_response(
        self,
        response: str,
        fallback_heading: str,
        fallback_content: str
    ) -> Dict[str, Any]:
        """LLM JSON yanıtını parse et"""
        
        # JSON bloğunu bul
        try:
            # Markdown code block içindeyse çıkar
            json_match = re.search(r'```(?:json)?\s*([\s\S]*?)```', response)
            if json_match:
                json_str = json_match.group(1).strip()
            else:
                # Direkt JSON olabilir
                json_start = response.find('{')
                json_end = response.rfind('}')
                if json_start >= 0 and json_end > json_start:
                    json_str = response[json_start:json_end + 1]
                else:
                    raise ValueError("JSON bulunamadı")
            
            result = json.loads(json_str)
            
            return {
                "heading": result.get("heading", fallback_heading),
                "enhanced_text": result.get("enhanced_text", fallback_content),
                "change_type": result.get("change_type", "content_restructured"),
                "explanation": result.get("explanation", "LLM ile iyileştirildi.")
            }
            
        except (json.JSONDecodeError, ValueError) as e:
            log_warning(f"LLM JSON parse hatası: {e}", "enhancer")
            # Fallback: ham yanıtı kullan
            return {
                "heading": fallback_heading,
                "enhanced_text": response.strip() if len(response.strip()) > 20 else fallback_content,
                "change_type": "content_restructured",
                "explanation": "LLM yanıtı yapılandırılamadı, ham iyileştirme uygulandı."
            }
    
    # ─────────────────────────────────────────
    #  ADIM 4: DOCX Oluşturma
    # ─────────────────────────────────────────
    
    def _generate_enhanced_docx(
        self,
        sections: List[EnhancedSection],
        original_name: str,
        session_id: str,
        original_content: bytes = None,
        file_type: str = ""
    ) -> str:
        """
        İyileştirilmiş doküman oluştur.
        
        Strateji:
        - DOCX orijinalse: Orijinal DOCX şablon olarak açılır, metinler güncellenir.
          (Orijinal görseller zaten korunur — _apply_to_original_docx bunu yapar)
        - PDF orijinalse: fpdf2 ile doğrudan PDF oluşturulur + orijinal görseller eklenir.
        - Diğerleri (PPTX/TXT): Sıfırdan DOCX oluşturulur.
        
        v2.40.1: Orijinal dosyadan görsel çıkarma ve iyileştirilmiş çıktıya gömme.
        """
        
        # Orijinal dosyadan görselleri çıkar (tüm format'lar için ortak)
        original_images = []
        if original_content:
            try:
                from app.services.document_processors.image_extractor import ImageExtractor
                img_extractor = ImageExtractor()
                ext_for_extract = file_type if file_type.startswith(".") else f".{file_type.lower()}"
                original_images = img_extractor.extract(original_content, ext_for_extract)
                if original_images:
                    log_system_event("INFO", f"Enhanced çıktı için {len(original_images)} görsel çıkarıldı", "enhancer")
            except Exception as e:
                log_warning(f"Enhanced çıktı görsel çıkarma hatası: {e}", "enhancer")
        
        # PDF orijinalse: doğrudan PDF oluştur (docx ara adımı yok)
        if file_type.upper() in ("PDF", ".PDF"):
            try:
                pdf_path = self._create_fresh_pdf(sections, original_name, session_id, original_images, original_content=original_content)
                if pdf_path and os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 0:
                    _enhanced_files[session_id] = pdf_path
                    log_system_event("INFO", f"Enhanced PDF oluşturuldu: {pdf_path}", "enhancer")
                    return pdf_path
            except Exception as e:
                log_system_event("WARNING", f"PDF oluşturma başarısız, DOCX fallback: {e}", "enhancer")
                # Fallback: DOCX olarak devam et
        
        # XLSX/XLS orijinalse: orijinal XLSX korunur + iyileştirilmiş sheet eklenir
        if original_content and file_type.upper().replace(".", "") in ("XLSX", "XLS"):
            try:
                xlsx_path = self._apply_to_original_xlsx(
                    original_content, sections, session_id, file_type
                )
                if xlsx_path and os.path.exists(xlsx_path) and os.path.getsize(xlsx_path) > 0:
                    _enhanced_files[session_id] = xlsx_path
                    log_system_event("INFO", f"Enhanced XLSX oluşturuldu: {xlsx_path}", "enhancer")
                    return xlsx_path
            except Exception as e:
                log_system_event("WARNING", f"XLSX oluşturma başarısız, DOCX fallback: {e}", "enhancer")
        
        from docx import Document
        _ = Document  # suppress pyflakes: used below in _apply/_create methods
        
        # DOCX ise orijinal dokümanı şablon olarak aç (görseller zaten korunur)
        if original_content and file_type.upper() in ("DOCX", ".DOCX"):
            doc = self._apply_to_original_docx(original_content, sections)
        else:
            # PPTX/TXT vb. veya PDF fallback için sıfırdan DOCX oluştur + görseller ekle
            doc = self._create_fresh_docx(sections, original_name, original_images)
        
        # Geçici DOCX dosyasına kaydet
        temp_docx_path = tempfile.mktemp(suffix=".docx", prefix=f"enhanced_{session_id}_")
        doc.save(temp_docx_path)
        
        # Session registry'ye ekle
        _enhanced_files[session_id] = temp_docx_path
        
        log_system_event("INFO", f"Enhanced DOCX oluşturuldu: {temp_docx_path}", "enhancer")
        return temp_docx_path
    
    def _apply_to_original_xlsx(
        self,
        original_content: bytes,
        sections: List[EnhancedSection],
        session_id: str,
        file_type: str = ".xlsx"
    ) -> str:
        """
        Orijinal XLSX dosyasını koruyarak iyileştirilmiş metni yeni sheet olarak ekler.
        
        Strateji:
        - Orijinal tüm sheet'ler aynen korunur (veri, format, görseller dahil)
        - Her iyileştirilmiş bölüm için '[Enhanced] OriginalSheetAdı' adında yeni sheet eklenir
        - Yeni sheet'te iyileştirilmiş metin satır satır yazılır
        - no_change bölümler atlanır (yeni sheet oluşturulmaz)
        
        Returns:
            Geçici XLSX dosya yolu
        """
        import io
        from openpyxl import load_workbook
        from openpyxl.styles import Font, PatternFill, Alignment
        
        wb = load_workbook(io.BytesIO(original_content))
        
        # İyileştirilmiş bölümleri yeni sheet olarak ekle
        added_count = 0
        for section in sections:
            if section.change_type == "no_change":
                continue
            
            text = self._get_section_text(section)
            if not text or not text.strip():
                continue
            
            # Sheet adı oluştur (Excel max 31 karakter)
            base_name = (section.heading or f"Bolum_{section.section_index + 1}")[:20]
            sheet_name = f"[E] {base_name}"
            
            # Aynı isimli sheet varsa numara ekle
            counter = 1
            original_sheet_name = sheet_name
            while sheet_name in wb.sheetnames:
                sheet_name = f"{original_sheet_name[:27]}_{counter}"
                counter += 1
            
            ws = wb.create_sheet(title=sheet_name)
            
            # Başlık satırı
            header_fill = PatternFill(start_color="2E86AB", end_color="2E86AB", fill_type="solid")
            header_font = Font(name="Calibri", size=12, bold=True, color="FFFFFF")
            
            ws.cell(row=1, column=1, value=f"İyileştirilmiş: {section.heading or f'Bölüm {section.section_index + 1}'}")
            ws.cell(row=1, column=1).font = header_font
            ws.cell(row=1, column=1).fill = header_fill
            ws.cell(row=1, column=1).alignment = Alignment(horizontal="left")
            
            # İyileştirme tipi
            ws.cell(row=2, column=1, value=f"Değişiklik: {section.change_type}")
            ws.cell(row=2, column=1).font = Font(name="Calibri", size=10, italic=True, color="666666")
            
            if section.explanation:
                ws.cell(row=2, column=2, value=section.explanation)
                ws.cell(row=2, column=2).font = Font(name="Calibri", size=10, italic=True, color="666666")
            
            # İyileştirilmiş metin satırları
            lines = text.split("\n")
            data_font = Font(name="Calibri", size=11)
            
            for li, line in enumerate(lines, start=4):
                # Pipe ile ayrılmış satırları sütunlara böl
                if " | " in line or "," in line:
                    separator = " | " if " | " in line else ","
                    parts = line.split(separator)
                    for ci, part in enumerate(parts):
                        cell = ws.cell(row=li, column=ci + 1, value=part.strip())
                        cell.font = data_font
                else:
                    cell = ws.cell(row=li, column=1, value=line)
                    cell.font = data_font
            
            # Sütun genişliklerini ayarla
            for col in ws.columns:
                max_len = 0
                for cell in col:
                    if cell.value:
                        max_len = max(max_len, len(str(cell.value)))
                ws.column_dimensions[col[0].column_letter].width = min(max_len + 2, 50)
            
            added_count += 1
        
        if added_count == 0:
            log_system_event("INFO", "XLSX: Hiç iyileştirme uygulanmadı, orijinal dosya döndürülüyor", "enhancer")
        
        # Geçici XLSX dosyasına kaydet
        ext = ".xlsx"
        temp_path = tempfile.mktemp(suffix=ext, prefix=f"enhanced_{session_id}_")
        wb.save(temp_path)
        
        log_system_event(
            "INFO",
            f"Enhanced XLSX: {added_count} iyileştirilmiş sheet eklendi (orijinal korundu)",
            "enhancer"
        )
        return temp_path
    
    # ─────────────────────────────────────────
    #  Görsel Eşleştirme Yardımcıları (DRY)
    # ─────────────────────────────────────────
    
    @staticmethod
    def _get_section_text(section: 'EnhancedSection') -> str:
        """
        Section'dan render edilecek metni döndürür.
        
        Kural: change_type "no_change" ise orijinal metin,
        aksi halde iyileştirilmiş metin kullanılır.
        
        Returns:
            Boş string-safe metin (asla None dönmez)
        """
        if section.change_type != "no_change" and section.enhanced_text:
            return section.enhanced_text
        return section.original_text or ""
    
    def _map_images_to_sections(
        self,
        sections: List['EnhancedSection'],
        original_images: list
    ) -> Dict[int, List[tuple]]:
        """
        Orijinal görselleri section'lara eşleştirir ve bölüm içi
        paragraf-relative pozisyonlarını hesaplar.
        
        Eşleştirme stratejisi (öncelik sırasıyla):
          1. Heading bazlı — görselin context_heading'i section heading'iyle eşleştirilir
          2. chunk_index bazlı — görselin sayfa/chunk numarası section index'iyle eşleştirilir
          3. Son section'a atama — hiçbir eşleşme bulunamazsa son bölüme eklenir
        
        Pozisyon hesaplama:
          - paragraph_index ≥ 0 → section'ın satır sayısına göre relative pozisyon
          - paragraph_index = -1 → pozisyon bilinmiyor, section sonuna eklenir
        
        Args:
            sections: İyileştirilmiş bölüm listesi
            original_images: ExtractedImage nesneleri listesi
            
        Returns:
            Dict[section_index, List[(relative_para_pos, img_obj)]]
            Her section'daki görseller pozisyona göre sıralıdır.
        """
        section_image_map: Dict[int, List[tuple]] = {}
        
        if not original_images:
            return section_image_map
        
        for img_obj in original_images:
            heading = getattr(img_obj, "context_heading", "") or ""
            chunk_idx = getattr(img_obj, "context_chunk_index", 0)
            para_idx = getattr(img_obj, "paragraph_index", -1)
            
            # ── 1. Heading bazlı section eşleştirmesi ──
            best_section = None
            for sec in sections:
                sec_heading = sec.heading or ""
                if heading and (heading in sec_heading or sec_heading in heading):
                    best_section = sec
                    break
            
            # ── 2. chunk_index ile section eşleştirmesi (fallback) ──
            if best_section is None:
                if chunk_idx < len(sections):
                    best_section = sections[chunk_idx]
                elif sections:
                    best_section = sections[-1]
            
            if best_section is None:
                continue
            
            sec_idx = best_section.section_index
            
            # ── 3. Bölüm içi relative pozisyon hesaplama ──
            relative_pos = -1
            if para_idx >= 0:
                text = self._get_section_text(best_section)
                total_lines = len([l for l in text.split("\n") if l.strip()])
                
                if total_lines > 0:
                    # Global paragraph_index → section-local pozisyona dönüştür
                    relative_pos = min(para_idx, total_lines - 1)
                    
                    # chunk_index üzerinden önceki section'ların satır sayısını çıkar
                    if chunk_idx < len(sections):
                        sec_start_approx = sum(
                            len([l for l in self._get_section_text(s).split("\n") if l.strip()])
                            for s in sections[:chunk_idx]
                        )
                        relative_pos = max(0, para_idx - sec_start_approx)
                        relative_pos = min(relative_pos, total_lines - 1)
            
            section_image_map.setdefault(sec_idx, []).append((relative_pos, img_obj))
        
        # Her section'daki görselleri pozisyona göre sırala
        # (pozisyonsuz görseller → sonuç olarak sonuncu)
        for sec_idx in section_image_map:
            section_image_map[sec_idx].sort(
                key=lambda x: x[0] if x[0] >= 0 else 999999
            )
        
        return section_image_map
    
    @staticmethod
    def _organize_images_at_positions(
        sec_imgs: List[tuple],
        total_paragraphs: int
    ) -> Dict[int, list]:
        """
        Bir section'daki görselleri paragraf pozisyonuna göre dict'e organize eder.
        
        Rendering sırasında her paragraftan sonra bu dict kontrol edilerek
        ilgili görseller doğru konuma eklenir.
        
        Args:
            sec_imgs: [(relative_para_pos, img_obj), ...] — _map_images_to_sections çıktısı
            total_paragraphs: Section'daki toplam paragraf sayısı
            
        Returns:
            Dict[para_position, [img_obj, ...]]
            Key = paragraf index'i (0-based) → o paragraftan sonra eklenecek görseller
            Key = total_paragraphs → section sonuna eklenecek görseller
        """
        imgs_at_pos: Dict[int, list] = {}
        
        for rel_pos, img_obj in sec_imgs:
            if rel_pos < 0:
                # Pozisyon bilinmiyor → section sonuna eklenecek
                imgs_at_pos.setdefault(total_paragraphs, []).append(img_obj)
            else:
                # Görseli bu paragraftan sonra ekle (bounds-safe)
                safe_pos = min(rel_pos, max(0, total_paragraphs - 1))
                imgs_at_pos.setdefault(safe_pos, []).append(img_obj)
        
        return imgs_at_pos
    
    def _create_fresh_pdf(self, sections: List[EnhancedSection], original_name: str, session_id: str, original_images: list = None, original_content: bytes = None) -> str:
        """
        fpdf2 ile doğrudan PDF oluştur.
        Saf Python — Word/COM/internet gerektirmez.
        Türkçe karakter desteği için Windows Arial TTF kullanılır.
        Markdown syntax temizlenerek düzgün tipografi ile render edilir.
        
        v2.40.1: Orijinal dosyadan çıkarılan görseller bölüm heading'lerine 
        eşleştirilerek ilgili bölümün sonuna eklenir.
        """
        from fpdf import FPDF
        
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        
        # Türkçe karakter desteği: Arial TTF (Windows sistem fontu)
        font_name = "Helvetica"  # Fallback
        has_bold = True  # Helvetica built-in bold var
        font_dir = r"C:\Windows\Fonts"
        arial_regular = os.path.join(font_dir, "arial.ttf")
        arial_bold = os.path.join(font_dir, "arialbd.ttf")
        
        if os.path.exists(arial_regular):
            font_name = "ArialTR"
            pdf.add_font(font_name, "", arial_regular)
            if os.path.exists(arial_bold):
                pdf.add_font(font_name, "B", arial_bold)
            else:
                has_bold = False
        
        pdf.set_font(font_name, size=11)
        
        # v3.3.0 [C6]: Orijinal PDF'den font bilgisi çıkar (varsa)
        original_body_size = 11  # Varsayılan
        try:
            import fitz  # PyMuPDF
            if original_content:
                orig_doc = fitz.open(stream=original_content, filetype="pdf")
                font_sizes = []
                # İlk 5 sayfa (performans) üzerinden font boyutlarını topla
                for page_num in range(min(5, len(orig_doc))):
                    page = orig_doc[page_num]
                    blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE).get("blocks", [])
                    for block in blocks:
                        for line in block.get("lines", []):
                            for span in line.get("spans", []):
                                fs = span.get("size", 0)
                                text = span.get("text", "").strip()
                                if 8 <= fs <= 14 and len(text) > 20:
                                    font_sizes.append(round(fs, 1))
                orig_doc.close()
                if font_sizes:
                    # En sık kullanılan font boyutu = body text boyutu
                    from collections import Counter
                    original_body_size = Counter(font_sizes).most_common(1)[0][0]
                    log_system_event("INFO", f"[C6] Orijinal PDF body font size: {original_body_size}pt", "enhancer")
        except Exception as e:
            log_warning(f"[C6] PDF font tespiti başarısız, varsayılan kullanılıyor: {e}", "enhancer")
        
        # Tespit edilen font boyutunu body text için kullan
        pdf.set_font(font_name, size=original_body_size)
        
        # Orijinal görselleri section + paragraf pozisyonuna eşleştir
        section_image_map = self._map_images_to_sections(sections, original_images)
        
        def _clean_markdown(text: str) -> str:
            """Markdown formatting'i temizle."""
            # Bold/italic markers
            text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
            text = re.sub(r'\*(.+?)\*', r'\1', text)
            text = re.sub(r'__(.+?)__', r'\1', text)
            text = re.sub(r'_(.+?)_', r'\1', text)
            # Inline code
            text = re.sub(r'`(.+?)`', r'\1', text)
            return text.strip()
        
        def _render_line(line: str):
            """Tek satırı markdown tipine göre render et."""
            stripped = line.strip()
            if not stripped:
                return
            
            # Heading tespiti: ## veya ### ile başlıyorsa
            if stripped.startswith('### '):
                # Alt-alt başlık (küçük bold)
                heading_text = _clean_markdown(stripped[4:])
                if has_bold:
                    pdf.set_font(font_name, "B", original_body_size)
                else:
                    pdf.set_font(font_name, size=original_body_size)
                pdf.multi_cell(0, 6, heading_text)
                pdf.ln(1)
                pdf.set_font(font_name, size=original_body_size)
            elif stripped.startswith('## '):
                # Alt başlık (orta bold)
                heading_text = _clean_markdown(stripped[3:])
                if has_bold:
                    pdf.set_font(font_name, "B", original_body_size + 1)
                else:
                    pdf.set_font(font_name, size=original_body_size + 1)
                pdf.multi_cell(0, 7, heading_text)
                pdf.ln(2)
                pdf.set_font(font_name, size=original_body_size)
            elif stripped.startswith('# '):
                # Ana başlık
                heading_text = _clean_markdown(stripped[2:])
                if has_bold:
                    pdf.set_font(font_name, "B", original_body_size + 3)
                else:
                    pdf.set_font(font_name, size=original_body_size + 3)
                pdf.multi_cell(0, 8, heading_text)
                pdf.ln(3)
                pdf.set_font(font_name, size=original_body_size)
            elif stripped.startswith(('- ', '* ', '• ')):
                # Madde işareti
                bullet_text = _clean_markdown(stripped[2:])
                pdf.multi_cell(0, 6, f"  •  {bullet_text}")
                pdf.ln(1)
            elif re.match(r'^\d+[\.\\)]\s', stripped):
                # Numaralı liste
                clean = _clean_markdown(stripped)
                pdf.multi_cell(0, 6, f"  {clean}")
                pdf.ln(1)
            elif stripped.startswith('---') or stripped.startswith('==='):
                # Ayırıcı çizgi — boşluk bırak
                pdf.ln(3)
            else:
                # Normal paragraf
                clean = _clean_markdown(stripped)
                pdf.multi_cell(0, 6, clean)
                pdf.ln(2)
        
        def _add_image_to_pdf(img_obj):
            """Tek bir görseli PDF'e ekle."""
            try:
                from PIL import Image as PILImage
                img_data = getattr(img_obj, "image_data", None)
                if not img_data:
                    return
                
                # PIL ile aç ve geçici dosyaya kaydet (fpdf2 bytes desteklemez)
                pil_img = PILImage.open(io.BytesIO(img_data))
                if pil_img.mode in ("RGBA", "P"):
                    pil_img = pil_img.convert("RGB")
                
                # Boyut hesapla: sayfa genişliğine sığdır (max 170mm)
                img_w, img_h = pil_img.size
                max_width_mm = 170  # A4 page width - margins
                
                # DPI hesabı: varsayılan 96 DPI
                w_mm = (img_w / 96.0) * 25.4
                h_mm = (img_h / 96.0) * 25.4
                
                # Ölçekle
                if w_mm > max_width_mm:
                    scale = max_width_mm / w_mm
                    w_mm = max_width_mm
                    h_mm *= scale
                
                # Sayfa taşması kontrolü
                available_h = pdf.h - pdf.get_y() - pdf.b_margin - 5
                if h_mm > available_h:
                    pdf.add_page()
                
                # Geçici dosyaya kaydet
                temp_img_path = tempfile.mktemp(suffix=".png")
                pil_img.save(temp_img_path, format="PNG")
                
                try:
                    pdf.image(temp_img_path, x=20, w=w_mm, h=h_mm)
                    pdf.ln(4)
                finally:
                    # Geçici dosyayı sil
                    try:
                        os.remove(temp_img_path)
                    except OSError:
                        pass
                        
            except Exception as e:
                log_warning(f"PDF görsel ekleme hatası: {e}", "enhancer")
        
        for section in sections:
            heading_text = section.heading or f"Bölüm {section.section_index + 1}"
            # Heading'deki markdown prefix'lerini temizle
            heading_text = re.sub(r'^#{1,6}\s+', '', heading_text).strip()
            heading_text = _clean_markdown(heading_text)
            
            # Section başlığı
            if has_bold:
                pdf.set_font(font_name, "B", original_body_size + 3)
            else:
                pdf.set_font(font_name, size=original_body_size + 3)
            pdf.cell(0, 10, heading_text, new_x="LMARGIN", new_y="NEXT")
            pdf.ln(2)
            
            # İçerik — görselleri paragraflar arasına yerleştir
            text = self._get_section_text(section)
            lines = [l for l in (text or "").split("\n")]
            visible_lines = [l for l in lines if l.strip()]
            sec_imgs = section_image_map.get(section.section_index, [])
            
            # Görselleri satır pozisyonlarına organize et
            imgs_at_line = self._organize_images_at_positions(sec_imgs, len(visible_lines))
            
            # Satırları ve görselleri sıralı render et
            visible_idx = 0
            if text:
                pdf.set_font(font_name, size=original_body_size)
                for paragraph_text in lines:
                    _render_line(paragraph_text)
                    if paragraph_text.strip():
                        # Bu satırdan sonra eklenecek görseller var mı?
                        for img_obj in imgs_at_line.get(visible_idx, []):
                            pdf.ln(2)
                            _add_image_to_pdf(img_obj)
                        visible_idx += 1
            
            # Section sonuna eklenmesi gereken görseller (pozisyonsuz)
            for img_obj in imgs_at_line.get(len(visible_lines), []):
                pdf.ln(2)
                _add_image_to_pdf(img_obj)
            
            pdf.ln(4)
        
        # Geçici dosyaya kaydet
        temp_pdf_path = tempfile.mktemp(suffix=".pdf", prefix=f"enhanced_{session_id}_")
        pdf.output(temp_pdf_path)
        
        return temp_pdf_path
    
    def _apply_to_original_docx(self, original_content: bytes, sections: List[EnhancedSection]):
        """
        Orijinal DOCX'i açıp sadece iyileştirilmiş bölümlerin paragraflarını günceller.
        Stiller, tablolar, resimler, header/footer korunur.
        """
        from docx import Document
        
        doc = Document(io.BytesIO(original_content))
        
        # Değişiklik gereken section'ları filtrele
        changed_sections = {
            s.section_index: s for s in sections if s.change_type != "no_change"
        }
        
        if not changed_sections:
            return doc  # Değişiklik yok, orijinali aynen döndür
        
        # Paragrafları heading bazlı section'lara eşle
        current_heading = "Giriş"
        section_idx = 0
        section_paragraphs: Dict[int, List] = {}  # section_idx → [paragraph_refs]
        section_headings: Dict[int, str] = {}  # section_idx → heading_text
        
        for para in doc.paragraphs:
            style_name = para.style.name if para.style else ""
            
            if style_name.startswith("Heading") and para.text.strip():
                # Önceki section'ı kapat
                if section_idx > 0 or section_paragraphs.get(0):
                    section_idx += 1 if section_paragraphs.get(section_idx) else 0
                
                current_heading = para.text.strip()
                section_headings[section_idx] = current_heading
                
                # Heading paragrafı eğer iyileştirilmiş heading varsa güncelle
                if section_idx in changed_sections:
                    new_heading = changed_sections[section_idx].heading
                    if new_heading and new_heading != current_heading:
                        self._update_paragraph_text(para, new_heading)
                
                # Bir sonraki section için hazırla
                section_idx += 1
                section_paragraphs.setdefault(section_idx, [])
            else:
                section_paragraphs.setdefault(section_idx, []).append(para)
        
        # Her değişen section için paragrafları güncelle
        for sec_idx, enhanced in changed_sections.items():
            paragraphs = section_paragraphs.get(sec_idx, [])
            if not paragraphs:
                continue
            
            # İyileştirilmiş metni satırlara böl
            new_lines = [line for line in enhanced.enhanced_text.split("\n") if line.strip()]
            
            # Mevcut paragrafları güncelle veya ekle
            for i, para in enumerate(paragraphs):
                if i < len(new_lines):
                    # Paragraf var → metnini güncelle (stil korunur)
                    self._update_paragraph_text(para, new_lines[i])
                else:
                    # Fazla paragrafları temizle (boşalt ama silme — stil korunsun)
                    self._update_paragraph_text(para, "")
            
            # Yeni metin paragraflardan fazlaysa kalan satırları son paragrafın
            # sonuna ekle (doküman yapısını bozmamak için)
            if len(new_lines) > len(paragraphs) and paragraphs:
                last_para = paragraphs[-1]
                remaining = new_lines[len(paragraphs):]
                current_text = last_para.text
                combined = current_text + "\n" + "\n".join(remaining) if current_text else "\n".join(remaining)
                self._update_paragraph_text(last_para, combined)
        
        return doc
    
    def _update_paragraph_text(self, para, new_text: str):
        """
        Paragrafın metnini güncellerken orijinal stil ve formatı korur.
        İlk metin run'ının formatını referans alarak metin run'larını günceller.
        
        KRİTİK: Inline görsel (drawing) içeren run'lar ASLA temizlenmez —
        bu sayede orijinal DOCX'teki görseller korunur.
        """
        if not para.runs:
            # Run yoksa direkt text güncelle
            para.text = new_text
            return
        
        # Run'ları iki gruba ayır: metin run'ları ve görsel (drawing) run'ları
        # OOXML namespace'leri: w:drawing ve wp:inline
        WML_DRAWING = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing'
        WP_INLINE = '{http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing}inline'
        
        text_runs = []
        image_runs = []
        
        for run in para.runs:
            has_drawing = (
                run._element.findall(f'.//{WML_DRAWING}') or
                run._element.findall(f'.//{WP_INLINE}')
            )
            if has_drawing:
                image_runs.append(run)
            else:
                text_runs.append(run)
        
        # Eğer paragrafta HİÇ metin run'ı yoksa (paragraf sadece görsel ise)
        # → hiçbir şey yapma, görseli koru
        if not text_runs:
            return
        
        # İlk metin run'ından font özelliklerini sakla
        first_run = text_runs[0]
        font_props = {}
        for prop in ('bold', 'italic', 'underline', 'size', 'name'):
            val = getattr(first_run.font, prop, None)
            if val is not None:
                font_props[prop] = val
        
        # Color'u da sakla
        try:
            if first_run.font.color and first_run.font.color.rgb:
                font_props['color_rgb'] = first_run.font.color.rgb
        except Exception as e:
            log_warning(f"Font color okuma hatası: {e}", "enhancer")
        
        # Sadece METIN run'larını temizle (görsel run'larına DOKUNMA)
        for run in text_runs:
            run.clear()
        
        # İlk metin run'ının metnini güncelle
        first_run.text = new_text
        
        # Saklanan font özelliklerini geri yükle
        for prop, val in font_props.items():
            if prop == 'color_rgb':
                try:
                    first_run.font.color.rgb = val
                except Exception as e:
                    log_warning(f"Font color geri yükleme hatası: {e}", "enhancer")
            else:
                try:
                    setattr(first_run.font, prop, val)
                except Exception as e:
                    log_warning(f"Font özelliği geri yükleme hatası ({prop}): {e}", "enhancer")
    
    def _create_fresh_docx(self, sections: List[EnhancedSection], original_name: str, original_images: list = None):
        """
        PDF/PPTX/TXT gibi DOCX olmayan dosyalar için sıfırdan DOCX oluşturur.
        Orijinal yapıyı korur — LLM meta-notları veya açıklama eklenmez.
        
        v2.40.2: Görseller bölüm sonuna topluca değil, paragraflar arasına
        orijinal sırasına göre (paragraph_index bazlı) yerleştirilir.
        """
        from docx import Document
        from docx.shared import Pt
        
        doc = Document()
        
        # Görselleri bölümlere ve bölüm içi pozisyonlara eşleştir
        section_image_map = self._map_images_to_sections(sections, original_images)
        
        # Dökümanı oluştur
        for section in sections:
            heading_text = section.heading or f"Bölüm {section.section_index + 1}"
            doc.add_heading(heading_text, level=1)
            
            text = self._get_section_text(section)
            paragraphs = [p.strip() for p in text.split("\n") if p.strip()]
            sec_imgs = section_image_map.get(section.section_index, [])
            
            # Görselleri paragraf pozisyonlarına organize et
            imgs_at_pos = self._organize_images_at_positions(sec_imgs, len(paragraphs))
            
            # Paragrafları ve görselleri sırayla ekle
            for p_idx, paragraph_text in enumerate(paragraphs):
                para = doc.add_paragraph(paragraph_text)
                para.style.font.size = Pt(11)
                
                # Bu paragraftan sonra eklenecek görseller var mı?
                for img_obj in imgs_at_pos.get(p_idx, []):
                    self._add_image_to_docx(doc, img_obj)
            
            # Section sonuna eklenmesi gereken görseller (pozisyonsuz)
            for img_obj in imgs_at_pos.get(len(paragraphs), []):
                self._add_image_to_docx(doc, img_obj)
        
        return doc
    
    def _add_image_to_docx(self, doc, img_obj):
        """Tek bir görseli DOCX'e ekler — boyut ölçekleme ve hata koruması dahil."""
        from docx.shared import Inches
        try:
            img_data = getattr(img_obj, "image_data", None)
            if not img_data:
                return
            
            img_stream = io.BytesIO(img_data)
            
            # Boyut hesapla: max 5.5 inch genişlik (A4 margins dahil)
            width = getattr(img_obj, "width", 0) or 400
            max_width_inch = 5.5
            w_inch = min(width / 96.0, max_width_inch)
            
            doc.add_picture(img_stream, width=Inches(w_inch))
        except Exception as e:
            log_warning(f"DOCX görsel ekleme hatası: {e}", "enhancer")
    
    def generate_selective_docx(
        self,
        original_content: bytes,
        sections: List[EnhancedSection],
        approved_indexes: List[int],
        session_id: str,
        file_type: str = ""
    ) -> str:
        """
        Sadece onaylanan section'ları uygulayarak yeni DOCX oluştur.
        Onaylanmayan bölümler orijinal haliyle kalır.
        
        Args:
            original_content: Orijinal dosya binary içeriği
            sections: Tüm EnhancedSection listesi
            approved_indexes: Kullanıcının onayladığı section index'leri
            session_id: Session kimliği
            file_type: Dosya tipi
        
        Returns:
            Geçici dosya yolu
        """
        # Sadece onaylanan section'ları aktif yap, diğerlerini no_change'e çevir
        selective_sections = []
        for s in sections:
            if s.section_index in approved_indexes:
                selective_sections.append(s)
            else:
                # Onaylanmayan → orijinal metin korunsun
                selective_sections.append(EnhancedSection(
                    section_index=s.section_index,
                    heading=s.heading,
                    original_text=s.original_text,
                    enhanced_text=s.original_text,
                    change_type="no_change",
                    explanation="Kullanıcı tarafından reddedildi.",
                    priority=s.priority,
                    violations=s.violations
                ))
        
        return self._generate_enhanced_docx(
            selective_sections, "", session_id,
            original_content=original_content, file_type=file_type
        )
    
    def to_dict(self, result: EnhancementResult) -> Dict[str, Any]:
        """EnhancementResult'ı JSON-serializable dict'e çevir"""
        return {
            "file_name": result.file_name,
            "file_type": result.file_type,
            "total_sections": result.total_sections,
            "enhanced_count": result.enhanced_count,
            "session_id": result.session_id,
            "error": result.error,
            "catboost_summary": result.catboost_summary,
            "sections": [
                {
                    "section_index": s.section_index,
                    "heading": s.heading,
                    "original_text": s.original_text,
                    "enhanced_text": s.enhanced_text,
                    "change_type": s.change_type,
                    "explanation": s.explanation,
                    "priority": s.priority,
                    "violations": s.violations,
                    "integrity_score": s.integrity_score,
                    "integrity_issues": s.integrity_issues,
                }
                for s in result.sections
            ]
        }


# ============================================
# Module-level helpers
# ============================================

def get_enhanced_file_path(session_id: str) -> Optional[str]:
    """Session ID ile geçici dosya yolunu getir"""
    return _enhanced_files.get(session_id)


def cleanup_enhanced_file(session_id: str):
    """Geçici dosyayı temizle"""
    import os
    path = _enhanced_files.pop(session_id, None)
    if path and os.path.exists(path):
        try:
            os.remove(path)
        except Exception as e:
            log_warning(f"Geçici dosya silme hatası ({path}): {e}", "enhancer")
