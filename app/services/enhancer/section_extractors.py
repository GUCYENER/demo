"""
VYRA L1 Support API - Section Extractors
==========================================
Dokümanları bölümlere ayırma (PDF, DOCX, XLSX, CSV, PPTX, TXT).

Her format için dedicated extractor metodu bulunur.
Heading detection + paragraf aralığı takibi dahildir.

Author: VYRA AI Team
Version: 1.0.0 (v3.3.3)
"""

import io
import re
from typing import Dict, Any, List, BinaryIO

from app.services.logging_service import log_error


class SectionExtractor:
    """
    Dokümanları bölümlere ayırma servisi.

    Desteklenen formatlar: PDF, DOCX, XLSX, CSV, PPTX, TXT

    Kullanım:
        extractor = SectionExtractor()
        sections = extractor.extract_sections(file_content, file_name, file_type)
    """

    def extract_sections(
        self,
        file_content: bytes,
        file_name: str,
        file_type: str
    ) -> List[Dict[str, Any]]:
        """Dokümanı bölümlere ayırır (heading bazlı veya sayfa bazlı)"""

        file_obj = io.BytesIO(file_content)
        sections = []

        try:
            if file_type == "PDF":
                sections = self._extract_pdf_sections(file_obj)
            elif file_type == "DOCX":
                sections = self._extract_docx_sections(file_obj)
            elif file_type == "XLSX":
                sections = self._extract_xlsx_sections(file_obj)
            elif file_type == "PPTX":
                sections = self._extract_pptx_sections(file_obj)
            elif file_type == "CSV":
                sections = self._extract_csv_sections(file_obj)
            elif file_type == "TXT":
                sections = self._extract_txt_sections(file_obj)
            else:
                # Bilinmeyen format — düz metin olarak dene
                file_obj.seek(0)
                text = file_obj.read().decode("utf-8", errors="replace")
                sections = [{"heading": "Genel", "content": text, "index": 0}]
        except Exception as e:
            log_error(f"Section extraction hatası ({file_type}): {e}", "enhancer")

        return sections

    def _extract_pdf_sections(self, file_obj: BinaryIO) -> List[Dict[str, Any]]:
        """PDF bölümlerini çıkar"""
        try:
            from pypdf import PdfReader
        except ImportError:
            from PyPDF2 import PdfReader

        file_obj.seek(0)
        reader = PdfReader(file_obj)

        all_text = ""
        for page in reader.pages:
            text = page.extract_text() or ""
            all_text += text + "\n"

        return self._split_text_by_headings(all_text)

    def _extract_docx_sections(self, file_obj: BinaryIO) -> List[Dict[str, Any]]:
        """DOCX bölümlerini çıkar (Word stilleri ile) — paragraf aralığı da kaydedilir"""
        from docx import Document

        file_obj.seek(0)
        doc = Document(file_obj)

        sections = []
        current_heading = "Giriş"
        current_content = []
        idx = 0
        para_start = 0  # Bu bölümün başladığı paragraf index'i

        for para_idx, para in enumerate(doc.paragraphs):
            style_name = para.style.name if para.style else ""
            text = para.text.strip()

            if not text:
                current_content.append("")
                continue

            if style_name.startswith("Heading"):
                # Önceki bölümü kaydet
                if current_content:
                    sections.append({
                        "heading": current_heading,
                        "content": "\n".join(current_content),
                        "index": idx,
                        "para_start": para_start,
                        "para_end": para_idx - 1
                    })
                    idx += 1
                current_heading = text
                current_content = []
                para_start = para_idx + 1  # Heading'den sonraki paragraf başlangıcı
            else:
                current_content.append(text)

        # Son bölümü kaydet
        if current_content:
            sections.append({
                "heading": current_heading,
                "content": "\n".join(current_content),
                "index": idx,
                "para_start": para_start,
                "para_end": len(doc.paragraphs) - 1
            })

        if not sections:
            # Heading yoksa tüm içeriği tek bölüm olarak al
            full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            sections = [{
                "heading": "Genel", "content": full_text, "index": 0,
                "para_start": 0, "para_end": len(doc.paragraphs) - 1
            }]

        return sections

    def _extract_xlsx_sections(self, file_obj: BinaryIO) -> List[Dict[str, Any]]:
        """
        XLSX bölümlerini çıkar — v3.3.0: Veri bloğu bazlı bölümleme.
        Büyük sheet'ler boş satır gap'lerine göre alt-section'lara ayrılır.
        Header satırı her alt-section'a prefix olarak eklenir.
        """
        from openpyxl import load_workbook

        file_obj.seek(0)
        wb = load_workbook(file_obj, data_only=True)
        sections = []
        global_idx = 0
        global_row_counter = 0

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]

            # Tüm satırları oku
            all_rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(cell) if cell is not None else "" for cell in row]
                has_data = any(c.strip() for c in cells)
                all_rows.append((cells, has_data))

            if not all_rows:
                continue

            # Header satırını tespit et (ilk veri satırı, kısa metin, benzersiz değerler)
            header_text = ""
            header_row_idx = -1
            for ri, (cells, has_data) in enumerate(all_rows):
                if has_data:
                    non_empty = [c for c in cells if c.strip()]
                    all_short = all(len(c) < 50 for c in non_empty)
                    all_unique = len(set(non_empty)) == len(non_empty)
                    if non_empty and all_short and all_unique and len(non_empty) >= 2:
                        header_text = " | ".join(non_empty)
                        header_row_idx = ri
                    break

            # Veri satırlarını boş satır gap'lerine göre data bloklarına böl
            data_blocks = []
            current_block = []
            gap_count = 0

            start_row = header_row_idx + 1 if header_row_idx >= 0 else 0

            for ri in range(start_row, len(all_rows)):
                cells, has_data = all_rows[ri]
                if has_data:
                    if gap_count >= 2 and current_block:
                        # 2+ boş satır = yeni data bloğu
                        data_blocks.append(current_block)
                        current_block = []
                    current_block.append(" | ".join(cells))
                    gap_count = 0
                else:
                    gap_count += 1

            if current_block:
                data_blocks.append(current_block)

            # Her data bloğunu section olarak ekle
            if not data_blocks:
                continue

            # Tek blok varsa sheet adıyla section oluştur
            if len(data_blocks) == 1:
                block_text = "\n".join(data_blocks[0])
                if header_text:
                    block_text = f"[Başlıklar: {header_text}]\n{block_text}"

                para_start = global_row_counter
                para_end = global_row_counter + len(data_blocks[0]) - 1
                sections.append({
                    "heading": sheet_name,
                    "content": block_text,
                    "index": global_idx,
                    "para_start": para_start,
                    "para_end": para_end
                })
                global_idx += 1
                global_row_counter += len(data_blocks[0])
            else:
                # Birden fazla blok — her bloğu alt-section yap
                for bi, block in enumerate(data_blocks):
                    block_text = "\n".join(block)
                    if header_text:
                        block_text = f"[Başlıklar: {header_text}]\n{block_text}"

                    block_heading = f"{sheet_name} — Bölüm {bi + 1}"
                    para_start = global_row_counter
                    para_end = global_row_counter + len(block) - 1
                    sections.append({
                        "heading": block_heading,
                        "content": block_text,
                        "index": global_idx,
                        "para_start": para_start,
                        "para_end": para_end
                    })
                    global_idx += 1
                    global_row_counter += len(block)

        return sections

    def _extract_csv_sections(self, file_obj: BinaryIO) -> List[Dict[str, Any]]:
        """
        CSV bölümlerini çıkar — v3.3.0: Satır bazlı bölümleme.
        CSVProcessor mantığıyla tutarlı: delimiter tespiti, header detection.
        """
        import csv as _csv

        file_obj.seek(0)
        raw = file_obj.read()

        # Encoding tespiti
        text = None
        try:
            from charset_normalizer import from_bytes
            result = from_bytes(raw).best()
            if result and result.encoding:
                text = str(result)
        except (ImportError, Exception):
            pass

        if text is None:
            for enc in ['utf-8', 'utf-8-sig', 'cp1254', 'iso-8859-9', 'latin-1']:
                try:
                    text = raw.decode(enc)
                    break
                except (UnicodeDecodeError, UnicodeError):
                    continue
            if text is None:
                text = raw.decode('utf-8', errors='replace')

        if not text or not text.strip():
            return []

        # Delimiter tespiti
        try:
            sample = "\n".join(text.split("\n")[:10])
            dialect = _csv.Sniffer().sniff(sample, delimiters=',;\t|')
            delimiter = dialect.delimiter
        except _csv.Error:
            counts = {',': text.count(','), ';': text.count(';'), '\t': text.count('\t')}
            delimiter = max(counts, key=counts.get)

        # CSV parse
        reader = _csv.reader(io.StringIO(text), delimiter=delimiter)
        all_rows = list(reader)

        if not all_rows:
            return []

        # Header tespiti
        has_header = False
        if len(all_rows) >= 2:
            first = all_rows[0]
            non_empty = [c for c in first if c.strip()]
            if non_empty and len(non_empty) >= 2:
                all_text_check = all(not c.strip().replace('.', '').replace(',', '').isdigit() for c in non_empty)
                all_short = all(len(c) < 50 for c in non_empty)
                all_unique = len(set(c.strip().lower() for c in non_empty)) == len(non_empty)
                has_header = all_text_check and all_short and all_unique

        header_row = all_rows[0] if has_header else None
        data_rows = all_rows[1:] if has_header else all_rows
        header_text = " | ".join(c.strip() for c in header_row if c.strip()) if header_row else ""

        # Bölümlere ayır (50 satır/bölüm)
        MAX_ROWS_PER_SECTION = 50
        sections = []

        for i in range(0, len(data_rows), MAX_ROWS_PER_SECTION):
            batch = data_rows[i:i + MAX_ROWS_PER_SECTION]
            row_texts = []
            for row in batch:
                row_text = " | ".join(c.strip() for c in row if c.strip())
                if row_text:
                    row_texts.append(row_text)

            if not row_texts:
                continue

            block_text = ""
            if header_text:
                block_text = f"[Başlıklar: {header_text}]\n"
            block_text += "\n".join(row_texts)

            section_heading = header_text or f"CSV Veri Bloğu {len(sections) + 1}"
            if len(sections) > 0 or i > 0:
                section_heading = f"{section_heading} — Bölüm {len(sections) + 1}"

            sections.append({
                "heading": section_heading,
                "content": block_text,
                "index": len(sections),
                "para_start": i,
                "para_end": i + len(row_texts) - 1
            })

        return sections

    def _extract_pptx_sections(self, file_obj: BinaryIO) -> List[Dict[str, Any]]:
        """PPTX her slaytı bir bölüm olarak çıkar — paragraf aralığı da kaydedilir"""
        from pptx import Presentation

        file_obj.seek(0)
        prs = Presentation(file_obj)
        sections = []
        global_para_counter = 0

        for idx, slide in enumerate(prs.slides):
            texts = []
            slide_title = f"Slayt {idx + 1}"

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if shape.shape_type is not None and shape.shape_type == 13:
                        continue  # Picture
                    if shape == slide.shapes.title:
                        slide_title = shape.text.strip()
                    else:
                        texts.append(shape.text.strip())

            if texts:
                para_start = global_para_counter
                para_end = global_para_counter + len(texts) - 1
                sections.append({
                    "heading": slide_title,
                    "content": "\n".join(texts),
                    "index": idx,
                    "para_start": para_start,
                    "para_end": para_end
                })
                global_para_counter += len(texts)

        return sections

    def _extract_txt_sections(self, file_obj: BinaryIO) -> List[Dict[str, Any]]:
        """TXT dosyasını satır bazlı bölümlere ayır"""
        file_obj.seek(0)
        text = file_obj.read().decode("utf-8", errors="replace")
        return self._split_text_by_headings(text)

    def _split_text_by_headings(self, text: str) -> List[Dict[str, Any]]:
        """Metni heading pattern'lerine göre bölümlere ayır — paragraf aralığı da kaydedilir"""
        lines = text.split("\n")
        sections = []
        current_heading = "Giriş"
        current_content = []
        idx = 0
        line_start = 0  # Bu bölümün başladığı satır index'i

        heading_patterns = [
            r'^\d+[\.\\)]\s+\S',
            r'^[A-ZÇĞİÖŞÜ][A-ZÇĞİÖŞÜ\s]{3,}$',
            r'^(?:BÖLÜM|MADDE|KISIM|BAŞLIK)\s',
        ]

        for line_idx, line in enumerate(lines):
            stripped = line.strip()
            if not stripped:
                current_content.append("")
                continue

            is_heading = False
            if len(stripped) < 80:
                for pattern in heading_patterns:
                    if re.match(pattern, stripped):
                        is_heading = True
                        break

            if is_heading:
                if current_content:
                    sections.append({
                        "heading": current_heading,
                        "content": "\n".join(current_content),
                        "index": idx,
                        "para_start": line_start,
                        "para_end": line_idx - 1
                    })
                    idx += 1
                current_heading = stripped
                current_content = []
                line_start = line_idx + 1  # Heading'den sonraki satır
            else:
                current_content.append(stripped)

        if current_content:
            sections.append({
                "heading": current_heading,
                "content": "\n".join(current_content),
                "index": idx,
                "para_start": line_start,
                "para_end": len(lines) - 1
            })

        # Hiç bölüm bulunmadıysa tüm metni tek bölüm yap
        if not sections:
            sections = [{
                "heading": "Genel", "content": text.strip(), "index": 0,
                "para_start": 0, "para_end": len(lines) - 1
            }]

        return sections
