"""
Document Maturity Analyzer
===========================
Dosya türüne göre RAG uyumluluğunu analiz eder.
Best practices kurallarına göre kategorik skorlama yapar.

@version 1.0.0
"""


import logging
import re
from collections import Counter
from typing import Dict, Any, List, BinaryIO

logger = logging.getLogger("vyra")

# v3.3.0 [B3]: Maturity analiz sonuçları cache'i (file hash bazlı, 10 dk TTL)
_maturity_cache: dict = {}


class MaturityRule:
    """Tek bir olgunluk kuralı"""
    def __init__(self, name: str, category: str, description: str, weight: float = 1.0):
        self.name = name
        self.category = category
        self.description = description
        self.weight = weight
        self.score = 100  # Varsayılan: tam puan
        self.status = "pass"  # pass / warning / fail
        self.detail = ""
        self.recommendation = ""  # v3.3.0 [B4]: Violation için çözüm önerisi
    
    def to_dict(self):
        result = {
            "name": self.name,
            "category": self.category,
            "description": self.description,
            "score": self.score,
            "status": self.status,
            "detail": self.detail
        }
        if self.recommendation:
            result["recommendation"] = self.recommendation
        return result


def analyze_pdf(file_obj: BinaryIO, file_name: str) -> Dict[str, Any]:
    """PDF dosyasını analiz eder"""
    try:
        from pypdf import PdfReader
    except ImportError:
        logger.debug("[MaturityAnalyzer] pypdf bulunamadı, PyPDF2 deneniyor", exc_info=True)
        from PyPDF2 import PdfReader
    
    rules: List[MaturityRule] = []
    
    file_obj.seek(0)
    reader = PdfReader(file_obj)
    total_pages = len(reader.pages)
    
    # Sayfa metinlerini çıkar
    page_texts = []
    all_text = ""
    for page in reader.pages:
        text = page.extract_text() or ""
        page_texts.append(text)
        all_text += text + "\n"
    
    lines = [l for l in all_text.split("\n") if l.strip()]
    
    # ─── KURAL 1: Başlık Hiyerarşisi ───
    rule = MaturityRule("Başlık Hiyerarşisi", "Yapı", "Heading 1/2/3 gibi başlık seviyeleri kullanılmalı")
    heading_patterns = [
        r'^\d+[\.\)]\s+\S',            # 1. veya 1) ile başlayan
        r'^\d+\.\d+[\.\)]\s+\S',       # 1.1 veya 1.1) ile başlayan
        r'^[A-ZÇĞİÖŞÜ][A-ZÇĞİÖŞÜ\s]{3,}$',  # TAMAMEN BÜYÜK HARF
        r'^[A-ZÇĞİÖŞÜ][\w\s]{2,50}$',          # Büyük harfle başlayan kısa satır
        r'^(?:BÖLÜM|MADDE|KISIM|BAŞLIK)\s',       # Türkçe bölüm kelimeleri
    ]
    heading_count = 0
    for line in lines:
        line_stripped = line.strip()
        if len(line_stripped) > 80:
            continue
        for pattern in heading_patterns:
            if re.match(pattern, line_stripped):
                heading_count += 1
                break
    
    if heading_count == 0:
        rule.score = 20
        rule.status = "fail"
        rule.detail = "Hiç başlık yapısı tespit edilemedi. Heading stilleri kullanılmalı."
    elif heading_count < 3:
        rule.score = 50
        rule.status = "warning"
        rule.detail = f"Sadece {heading_count} başlık tespit edildi. Daha fazla başlık hiyerarşisi önerilir."
    else:
        rule.score = 100
        rule.detail = f"{heading_count} başlık tespit edildi."
    rules.append(rule)
    
    # ─── KURAL 2: Tablo Formatı ───
    rule = MaturityRule("Tablo Formatı", "Tablo", "Tablolar düzenli satır/sütunlarla oluşturulmalı, merge hücre olmamalı")
    table_lines = [l for l in lines if '|' in l or '\t' in l]
    if table_lines:
        # Pipe ayrımı kontrol
        pipe_counts = [l.count('|') for l in table_lines if '|' in l]
        if pipe_counts:
            avg_pipes = sum(pipe_counts) / len(pipe_counts)
            variance = sum((c - avg_pipes)**2 for c in pipe_counts) / len(pipe_counts) if len(pipe_counts) > 1 else 0
            if variance > 4:
                rule.score = 60
                rule.status = "warning"
                rule.detail = "Tablo yapısında tutarsızlık tespit edildi. Sütun sayıları farklılık gösteriyor."
            else:
                rule.score = 100
                rule.detail = f"{len(table_lines)} tablo satırı tespit edildi, format tutarlı."
        else:
            rule.score = 85
            rule.detail = "Tab ayraçlı tablo tespit edildi."
    else:
        rule.score = 100  # Tablo yoksa sorun değil
        rule.detail = "Tablo içeriği tespit edilmedi (sorun değil)."
        rule.status = "pass"
    rules.append(rule)
    
    # ─── KURAL 3: Aranabilir Metin ───
    rule = MaturityRule("Aranabilir Metin", "Metin", "PDF metin tabanlı olmalı, taranmış görsel olmamalı")
    if total_pages > 0:
        avg_chars_per_page = len(all_text) / total_pages
        if avg_chars_per_page < 50:
            rule.score = 10
            rule.status = "fail"
            rule.detail = f"Sayfa başına ortalama {avg_chars_per_page:.0f} karakter. Muhtemelen taranmış (görsel) PDF. OCR gereklidir."
        elif avg_chars_per_page < 200:
            rule.score = 50
            rule.status = "warning"
            rule.detail = f"Sayfa başına ortalama {avg_chars_per_page:.0f} karakter. Metin içeriği düşük."
        else:
            rule.score = 100
            rule.detail = f"Sayfa başına ortalama {avg_chars_per_page:.0f} karakter. Yeterli metin içeriği var."
    rules.append(rule)
    
    # ─── KURAL 4: Metin vs Görsel Oranı ───
    rule = MaturityRule("Metin Yoğunluğu", "Metin", "Önemli bilgiler metin olarak yer almalı, görsel olmamalı")
    if total_pages > 0:
        chars_per_page = len(all_text.replace(" ", "").replace("\n", "")) / total_pages
        if chars_per_page < 100:
            rule.score = 30
            rule.status = "fail"
            rule.detail = "Çok düşük metin yoğunluğu. Bilgiler muhtemelen görsel içinde."
        elif chars_per_page < 300:
            rule.score = 65
            rule.status = "warning"
            rule.detail = "Metin yoğunluğu orta düzeyde."
        else:
            rule.score = 100
            rule.detail = "Metin yoğunluğu yeterli."
    rules.append(rule)
    
    # ─── KURAL 5: Gereksiz İçerik (header/footer tekrarı) ───
    rule = MaturityRule("Gereksiz İçerik", "Temizlik", "Tekrarlayan header/footer, sayfa numaraları temizlenmeli")
    if len(page_texts) > 2:
        first_lines = [t.split("\n")[0].strip() for t in page_texts if t.strip()]
        last_lines = [t.split("\n")[-1].strip() for t in page_texts if t.strip()]
        
        # İlk satır tekrarı (header)
        header_repeat = 0
        if first_lines:
            fl_counter = Counter(first_lines)
            most_common_fl = fl_counter.most_common(1)[0]
            if most_common_fl[1] > len(page_texts) * 0.5 and len(most_common_fl[0]) > 3:
                header_repeat = most_common_fl[1]
        
        # Son satır tekrarı (footer)
        footer_repeat = 0
        if last_lines:
            ll_counter = Counter(last_lines)
            most_common_ll = ll_counter.most_common(1)[0]
            if most_common_ll[1] > len(page_texts) * 0.5 and len(most_common_ll[0]) > 3:
                footer_repeat = most_common_ll[1]
        
        if header_repeat > 0 or footer_repeat > 0:
            rule.score = 60
            rule.status = "warning"
            parts = []
            if header_repeat: parts.append(f"header ({header_repeat} sayfada tekrar)")
            if footer_repeat: parts.append(f"footer ({footer_repeat} sayfada tekrar)")
            rule.detail = f"Tekrarlayan {', '.join(parts)} tespit edildi. Chunk kalitesini düşürebilir."
        else:
            rule.score = 100
            rule.detail = "Tekrarlayan header/footer tespit edilmedi."
    else:
        rule.score = 100
        rule.detail = "Sayfa sayısı çok az, kontrol atlandı."
    rules.append(rule)
    
    # ─── KURAL 6: Türkçe Karakter ───
    rule = MaturityRule("Türkçe Karakter", "Encoding", "ğ, ü, ş, ı, ö, ç karakterleri doğru encode edilmeli")
    turkish_chars = set("ğüşıöçĞÜŞİÖÇ")
    found_turkish = sum(1 for c in all_text if c in turkish_chars)
    
    # Bozuk karakter pattern'leri
    broken_patterns = [
        r'Ã¼', r'Ã§', r'Ã¶', r'Ä±', r'ÅŸ', r'Äž',  # UTF-8 bozuk
        r'\?{3,}',  # Üç veya daha fazla soru işareti (karakter kaybı)
    ]
    broken_count = 0
    for pattern in broken_patterns:
        broken_count += len(re.findall(pattern, all_text))
    
    if broken_count > 5:
        rule.score = 30
        rule.status = "fail"
        rule.detail = f"{broken_count} bozuk karakter encoding tespit edildi. Türkçe karakterler düzgün okumaz."
    elif broken_count > 0:
        rule.score = 70
        rule.status = "warning"
        rule.detail = f"{broken_count} olası encoding sorunu tespit edildi."
    elif found_turkish == 0 and len(all_text) > 500:
        rule.score = 80
        rule.status = "warning"
        rule.detail = "Hiç Türkçe karakter bulunamadı. Dosya Türkçe içerik içermiyor olabilir."
    else:
        rule.score = 100
        rule.detail = f"Türkçe karakter encoding sorunsuz ({found_turkish} Türkçe karakter bulundu)."
    rules.append(rule)
    
    return _build_result(rules, file_name, "PDF", total_pages)


def analyze_docx(file_obj: BinaryIO, file_name: str) -> Dict[str, Any]:
    """DOCX dosyasını analiz eder"""
    from docx import Document
    
    rules: List[MaturityRule] = []
    
    file_obj.seek(0)
    doc = Document(file_obj)
    
    paragraphs = doc.paragraphs
    total_paragraphs = len(paragraphs)
    
    # ─── KURAL 1: Word Stilleri ───
    rule = MaturityRule("Word Stilleri", "Yapı", "Heading 1/2/3 stilleri kullanılmalı, manuel kalın yazıdan kaçınılmalı")
    heading_styles = [p for p in paragraphs if p.style and p.style.name and p.style.name.startswith('Heading')]
    heading_count = len(heading_styles)
    
    if heading_count == 0:
        rule.score = 20
        rule.status = "fail"
        rule.detail = "Hiç Heading stili kullanılmamış. Word'ün Heading 1/2/3 stillerini kullanın."
    elif heading_count < 3:
        rule.score = 50
        rule.status = "warning"
        rule.detail = f"Sadece {heading_count} Heading stili tespit edildi. Daha fazla bölümleme önerilir."
    else:
        rule.score = 100
        rule.detail = f"{heading_count} Heading stili tespit edildi."
    rules.append(rule)
    
    # ─── KURAL 2: Metin Kutusu ───
    rule = MaturityRule("Metin Kutusu", "İçerik", "Text box içindeki metinler işlenemez, normal paragraf kullanılmalı")
    # python-docx ile inline shape sayısı
    try:
        inline_shapes = doc.inline_shapes
        _ = len(inline_shapes) if inline_shapes else 0
    except Exception:
        logger.debug("[MaturityAnalyzer] inline_shapes okuma hatası", exc_info=True)
        _ = 0
    
    # Floating shapes (text boxes) - XML tabanında kontrol
    try:
        textbox_count = 0
        for element in doc.element.body.iter():
            if element.tag.endswith('}txbxContent') or element.tag.endswith('}textbox'):
                textbox_count += 1
        if textbox_count > 0:
            rule.score = 40
            rule.status = "fail"
            rule.detail = f"{textbox_count} metin kutusu (text box) tespit edildi. İçindeki metinler işlenemez."
        else:
            rule.score = 100
            rule.detail = "Metin kutusu tespit edilmedi."
    except Exception:
        logger.debug("[MaturityAnalyzer] Metin kutusu kontrolü hatası", exc_info=True)
        rule.score = 90
        rule.detail = "Metin kutusu kontrolü yapılamadı (sorun değil)."
    rules.append(rule)
    
    # ─── KURAL 3: Tablo Başlık Satırı ───
    rule = MaturityRule("Tablo Başlık Satırı", "Tablo", "Tabloların ilk satırı başlık olmalı")
    tables = doc.tables
    if tables:
        tables_with_header = 0
        for table in tables:
            if table.rows:
                _ = " ".join(cell.text.strip() for cell in table.rows[0].cells)
                # Başlık heuristik: kısa, sayısal olmayan, farklı değerler
                cells = [cell.text.strip() for cell in table.rows[0].cells if cell.text.strip()]
                if cells and len(set(cells)) == len(cells) and all(len(c) < 50 for c in cells):
                    tables_with_header += 1
        
        ratio = tables_with_header / len(tables)
        if ratio >= 0.8:
            rule.score = 100
            rule.detail = f"{len(tables)} tablonun {tables_with_header} tanesinde başlık satırı tespit edildi."
        elif ratio >= 0.5:
            rule.score = 65
            rule.status = "warning"
            rule.detail = f"{len(tables)} tablonun sadece {tables_with_header} tanesinde başlık satırı var."
        else:
            rule.score = 35
            rule.status = "fail"
            rule.detail = f"{len(tables)} tablonun çoğunda başlık satırı eksik."
    else:
        rule.score = 100
        rule.detail = "Tablo içeriği yok (sorun değil)."
    rules.append(rule)
    
    # ─── KURAL 4: Görseller vs Metin ───
    rule = MaturityRule("Görseller vs Metin", "İçerik", "Önemli bilgiler metin olarak yazılmalı, görselde olmamalı")
    try:
        image_count = len(doc.inline_shapes) if doc.inline_shapes else 0
    except Exception:
        logger.debug("[MaturityAnalyzer] Görsel sayısı okuma hatası", exc_info=True)
        image_count = 0
    
    text_paragraphs = [p for p in paragraphs if p.text.strip()]
    if total_paragraphs > 0 and image_count > 0:
        image_ratio = image_count / max(len(text_paragraphs), 1)
        if image_ratio > 0.5:
            rule.score = 40
            rule.status = "fail"
            rule.detail = f"{image_count} görsel, {len(text_paragraphs)} metin paragrafı. Görsel oranı çok yüksek."
        elif image_ratio > 0.2:
            rule.score = 70
            rule.status = "warning"
            rule.detail = f"{image_count} görsel, {len(text_paragraphs)} metin paragrafı."
        else:
            rule.score = 100
            rule.detail = f"Metin/görsel oranı dengeli ({image_count} görsel)."
    else:
        rule.score = 100
        rule.detail = "Sorunsuz."
    rules.append(rule)
    
    # ─── KURAL 5: Liste Formatı ───
    rule = MaturityRule("Liste Formatı", "Yapı", "Madde işaretli ve numaralı listeler Word formatında olmalı")
    list_styles = [p for p in paragraphs if p.style and ('List' in (p.style.name or '') or 
                   p.paragraph_format.left_indent is not None and p.text.strip())]
    manual_list = [p for p in paragraphs if re.match(r'^[\-\•\*]\s', p.text.strip()) or 
                   re.match(r'^\d+[\.\)]\s', p.text.strip())]
    
    if manual_list and not list_styles:
        rule.score = 50
        rule.status = "warning"
        rule.detail = f"{len(manual_list)} manuel liste satırı tespit edildi. Word'ün liste stillerini kullanmanız önerilir."
    elif list_styles:
        rule.score = 100
        rule.detail = f"{len(list_styles)} Word liste stili kullanılmış."
    else:
        rule.score = 90
        rule.detail = "Liste içeriği tespit edilmedi."
    rules.append(rule)
    
    # ─── KURAL 6: Gereksiz Boşluklar ───
    rule = MaturityRule("Gereksiz Boşluklar", "Temizlik", "Fazla boş satırlar ve sayfa sonları temizlenmeli")
    empty_paragraphs = [p for p in paragraphs if not p.text.strip()]
    empty_ratio = len(empty_paragraphs) / max(total_paragraphs, 1)
    
    if empty_ratio > 0.4:
        rule.score = 35
        rule.status = "fail"
        rule.detail = f"Paragrafların %{empty_ratio*100:.0f}'i boş. Fazla boşluklar chunk kalitesini düşürür."
    elif empty_ratio > 0.2:
        rule.score = 65
        rule.status = "warning"
        rule.detail = f"Toplam {len(empty_paragraphs)} boş paragraf var (%{empty_ratio*100:.0f})."
    else:
        rule.score = 100
        rule.detail = "Boşluk oranı normal."
    rules.append(rule)
    
    return _build_result(rules, file_name, "DOCX", total_paragraphs)


def analyze_xlsx(file_obj: BinaryIO, file_name: str) -> Dict[str, Any]:
    """XLSX dosyasını analiz eder"""
    import openpyxl
    
    rules: List[MaturityRule] = []
    
    file_obj.seek(0)
    wb = openpyxl.load_workbook(file_obj, data_only=False)
    total_sheets = len(wb.sheetnames)
    
    # ─── KURAL 1: İlk Satır Başlık ───
    rule = MaturityRule("İlk Satır Başlık", "Yapı", "Her sayfanın ilk satırı sütun başlıklarını içermeli")
    sheets_with_header = 0
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        if ws.max_row and ws.max_row > 0:
            first_row = [cell.value for cell in ws[1] if cell.value is not None]
            if first_row:
                # Başlık heuristik: kısa metin, sayısal değil, benzersiz
                all_text_vals = all(isinstance(v, str) and len(str(v)) < 50 for v in first_row)
                unique = len(set(str(v) for v in first_row)) == len(first_row)
                if all_text_vals and unique and len(first_row) >= 2:
                    sheets_with_header += 1
    
    if total_sheets > 0:
        ratio = sheets_with_header / total_sheets
        if ratio >= 0.8:
            rule.score = 100
            rule.detail = f"{total_sheets} sayfanın {sheets_with_header} tanesinde başlık satırı var."
        elif ratio >= 0.5:
            rule.score = 60
            rule.status = "warning"
            rule.detail = f"{total_sheets} sayfanın sadece {sheets_with_header} tanesinde başlık satırı tespit edildi."
        else:
            rule.score = 25
            rule.status = "fail"
            rule.detail = f"Başlık satırı eksik. {total_sheets} sayfanın {sheets_with_header} tanesinde var."
    rules.append(rule)
    
    # ─── KURAL 2: Merge Hücreler ───
    rule = MaturityRule("Merge Hücreler", "Hücre", "Birleştirilmiş hücreler veri kaybına neden olabilir (processor otomatik çözümler)")
    total_merged = 0
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        total_merged += len(ws.merged_cells.ranges)
    
    if total_merged > 20:
        rule.score = 50
        rule.status = "warning"
        rule.detail = f"{total_merged} birleştirilmiş hücre aralığı var. Çoğu otomatik çözümlenir ama karmaşık merge yapıları sorun çıkarabilir."
    elif total_merged > 5:
        rule.score = 75
        rule.status = "warning"
        rule.detail = f"{total_merged} birleştirilmiş hücre var. Otomatik çözümleme uygulanır."
    elif total_merged > 0:
        rule.score = 90
        rule.detail = f"{total_merged} birleştirilmiş hücre var (otomatik çözümlenir)."
    else:
        rule.score = 100
        rule.detail = "Birleştirilmiş hücre yok."
    rules.append(rule)
    
    # ─── KURAL 3: Boş Satır/Sütun ───
    rule = MaturityRule("Boş Satır/Sütun", "Veri", "Veriler arasında boşluk bırakılmamalı")
    empty_row_gaps = 0
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        prev_empty = False
        data_started = False
        for row in ws.iter_rows(max_row=min(ws.max_row or 0, 200), values_only=True):
            has_data = any(cell is not None and str(cell).strip() for cell in row)
            if has_data:
                data_started = True
                if prev_empty and data_started:
                    empty_row_gaps += 1
                prev_empty = False
            else:
                if data_started:
                    prev_empty = True
    
    if empty_row_gaps > 5:
        rule.score = 35
        rule.status = "fail"
        rule.detail = f"{empty_row_gaps} yerde veri blokları arasında boşluk var. Bölümleme sorunlarına yol açar."
    elif empty_row_gaps > 2:
        rule.score = 65
        rule.status = "warning"
        rule.detail = f"{empty_row_gaps} yerde veri arasında boşluk tespit edildi."
    else:
        rule.score = 100
        rule.detail = "Veri blokları arasında boşluk yok."
    rules.append(rule)
    
    # ─── KURAL 4: Tutarlı Veri Tipi ───
    rule = MaturityRule("Tutarlı Veri Tipi", "Veri", "Her sütunda aynı veri tipi olmalı")
    type_issues = 0
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        if not ws.max_row or ws.max_row < 3:
            continue
        max_col = min(ws.max_column or 1, 20)
        for col_idx in range(1, max_col + 1):
            types_in_col = set()
            for row_idx in range(2, min(ws.max_row + 1, 50)):  # İlk 50 satır
                cell = ws.cell(row=row_idx, column=col_idx)
                if cell.value is not None:
                    types_in_col.add(type(cell.value).__name__)
            if len(types_in_col) > 2:  # 3+ farklı tip
                type_issues += 1
    
    if type_issues > 5:
        rule.score = 40
        rule.status = "fail"
        rule.detail = f"{type_issues} sütunda karışık veri tipleri tespit edildi."
    elif type_issues > 2:
        rule.score = 65
        rule.status = "warning"
        rule.detail = f"{type_issues} sütunda tip tutarsızlığı var."
    else:
        rule.score = 100
        rule.detail = "Veri tipleri tutarlı."
    rules.append(rule)
    
    # ─── KURAL 5: Açıklama Satırları ───
    rule = MaturityRule("Açıklama Satırları", "İçerik", "Veri üstüne açıklama satırları eklenmemeli")
    description_rows = 0
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        if not ws.max_row or ws.max_row < 2:
            continue
        # İlk 5 satırda tek hücrede çok uzun metin varsa açıklama satırı
        for row_idx in range(1, min(6, ws.max_row + 1)):
            row_values = [ws.cell(row=row_idx, column=c).value for c in range(1, min((ws.max_column or 1) + 1, 20))]
            non_empty = [v for v in row_values if v is not None and str(v).strip()]
            if len(non_empty) == 1 and len(str(non_empty[0])) > 80:
                description_rows += 1
    
    if description_rows > 3:
        rule.score = 40
        rule.status = "fail"
        rule.detail = f"{description_rows} açıklama satırı tespit edildi. Ayrı bir sayfada tutulmalı."
    elif description_rows > 0:
        rule.score = 70
        rule.status = "warning"
        rule.detail = f"{description_rows} olası açıklama satırı tespit edildi."
    else:
        rule.score = 100
        rule.detail = "Açıklama satırı tespit edilmedi."
    rules.append(rule)
    
    # ─── KURAL 6: Formül vs Değer ───
    rule = MaturityRule("Formül vs Değer", "İçerik", "Formüller işlenemez, 'Yapıştır > Değer' kullanılmalı")
    # openpyxl data_only=False ile formüller görülür
    formula_count = 0
    total_cells = 0
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        for row in ws.iter_rows(max_row=min(ws.max_row or 0, 100), max_col=min(ws.max_column or 1, 20)):
            for cell in row:
                if cell.value is not None:
                    total_cells += 1
                    if isinstance(cell.value, str) and cell.value.startswith('='):
                        formula_count += 1
    
    if total_cells > 0:
        formula_ratio = formula_count / total_cells
        if formula_ratio > 0.3:
            rule.score = 30
            rule.status = "fail"
            rule.detail = f"{formula_count} formüllü hücre tespit edildi (%{formula_ratio*100:.0f}). 'Değer Olarak Yapıştır' önerilir."
        elif formula_ratio > 0.1:
            rule.score = 60
            rule.status = "warning"
            rule.detail = f"{formula_count} formüllü hücre var (%{formula_ratio*100:.0f})."
        elif formula_count > 0:
            rule.score = 85
            rule.detail = f"Az sayıda formül ({formula_count} hücre)."
        else:
            rule.score = 100
            rule.detail = "Formül içeren hücre yok."
    else:
        rule.score = 100
        rule.detail = "Veri bulunamadı."
    rules.append(rule)
    
    # ─── KURAL 7: Gizli Sheet Kontrolü ───
    rule = MaturityRule("Gizli Sheet", "Yapı", "Gizli sayfalar chunk'a dahil edilir, istenmeyen veri sızabilir")
    hidden_sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        if ws.sheet_state != 'visible':
            hidden_sheets.append(sheet_name)
    
    if hidden_sheets:
        rule.score = 70
        rule.status = "warning"
        rule.detail = f"{len(hidden_sheets)} gizli sayfa tespit edildi ({', '.join(hidden_sheets[:3])}). Bu sayfalar da işlenecektir."
    else:
        rule.score = 100
        rule.detail = "Gizli sayfa yok."
    rules.append(rule)
    
    # ─── KURAL 8: Veri Boyutu Uyarısı ───
    rule = MaturityRule("Veri Boyutu", "Performans", "Çok büyük dosyalar işleme süresini uzatır")
    total_rows = 0
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        total_rows += ws.max_row or 0
    
    if total_rows > 50000:
        rule.score = 30
        rule.status = "fail"
        rule.detail = f"Toplam {total_rows:,} satır. Çok büyük dosya — işleme süresi uzun olabilir, bellek kullanımı yüksek."
    elif total_rows > 10000:
        rule.score = 65
        rule.status = "warning"
        rule.detail = f"Toplam {total_rows:,} satır. Büyük dosya — chunk sayısı yüksek olacaktır."
    else:
        rule.score = 100
        rule.detail = f"Toplam {total_rows:,} satır. Normal boyut."
    rules.append(rule)
    
    return _build_result(rules, file_name, "XLSX", total_sheets)


def analyze_xls(file_obj: BinaryIO, file_name: str) -> Dict[str, Any]:
    """XLS (.xls) dosyasını xlrd ile analiz eder"""
    import xlrd
    
    rules: List[MaturityRule] = []
    
    file_obj.seek(0)
    content = file_obj.read()
    wb = xlrd.open_workbook(file_contents=content)
    total_sheets = wb.nsheets
    
    # ─── KURAL 1: İlk Satır Başlık ───
    rule = MaturityRule("İlk Satır Başlık", "Yapı", "Her sayfanın ilk satırı sütun başlıklarını içermeli")
    sheets_with_header = 0
    for sheet_idx in range(wb.nsheets):
        sheet = wb.sheet_by_index(sheet_idx)
        if sheet.nrows > 0:
            first_row = [sheet.cell_value(0, c) for c in range(sheet.ncols) if sheet.cell_value(0, c)]
            if first_row:
                all_text_vals = all(isinstance(v, str) and len(str(v)) < 50 for v in first_row)
                unique = len(set(str(v) for v in first_row)) == len(first_row)
                if all_text_vals and unique and len(first_row) >= 2:
                    sheets_with_header += 1
    
    if total_sheets > 0:
        ratio = sheets_with_header / total_sheets
        if ratio >= 0.8:
            rule.score = 100
            rule.detail = f"{total_sheets} sayfanın {sheets_with_header} tanesinde başlık satırı var."
        elif ratio >= 0.5:
            rule.score = 60
            rule.status = "warning"
            rule.detail = f"{total_sheets} sayfanın sadece {sheets_with_header} tanesinde başlık satırı tespit edildi."
        else:
            rule.score = 25
            rule.status = "fail"
            rule.detail = f"Başlık satırı eksik. {total_sheets} sayfanın {sheets_with_header} tanesinde var."
    rules.append(rule)
    
    # ─── KURAL 2: Boş Satır/Sütun ───
    rule = MaturityRule("Boş Satır/Sütun", "Veri", "Veriler arasında boşluk bırakılmamalı")
    empty_row_gaps = 0
    for sheet_idx in range(wb.nsheets):
        sheet = wb.sheet_by_index(sheet_idx)
        prev_empty = False
        data_started = False
        for row_idx in range(min(sheet.nrows, 200)):
            row_values = [sheet.cell_value(row_idx, c) for c in range(sheet.ncols)]
            has_data = any(v is not None and str(v).strip() for v in row_values)
            if has_data:
                data_started = True
                if prev_empty and data_started:
                    empty_row_gaps += 1
                prev_empty = False
            else:
                if data_started:
                    prev_empty = True
    
    if empty_row_gaps > 5:
        rule.score = 35
        rule.status = "fail"
        rule.detail = f"{empty_row_gaps} yerde veri blokları arasında boşluk var."
    elif empty_row_gaps > 2:
        rule.score = 65
        rule.status = "warning"
        rule.detail = f"{empty_row_gaps} yerde veri arasında boşluk tespit edildi."
    else:
        rule.score = 100
        rule.detail = "Veri blokları arasında boşluk yok."
    rules.append(rule)
    
    # ─── KURAL 3: Tutarlı Veri Tipi ───
    rule = MaturityRule("Tutarlı Veri Tipi", "Veri", "Her sütunda aynı veri tipi olmalı")
    type_issues = 0
    for sheet_idx in range(wb.nsheets):
        sheet = wb.sheet_by_index(sheet_idx)
        if sheet.nrows < 3:
            continue
        max_col = min(sheet.ncols, 20)
        for col_idx in range(max_col):
            types_in_col = set()
            for row_idx in range(1, min(sheet.nrows, 50)):
                cell_type = sheet.cell_type(row_idx, col_idx)
                if cell_type != xlrd.XL_CELL_EMPTY:
                    types_in_col.add(cell_type)
            if len(types_in_col) > 2:
                type_issues += 1
    
    if type_issues > 5:
        rule.score = 40
        rule.status = "fail"
        rule.detail = f"{type_issues} sütunda karışık veri tipleri tespit edildi."
    elif type_issues > 2:
        rule.score = 65
        rule.status = "warning"
        rule.detail = f"{type_issues} sütunda tip tutarsızlığı var."
    else:
        rule.score = 100
        rule.detail = "Veri tipleri tutarlı."
    rules.append(rule)
    
    # ─── KURAL 4: Merge Hücreler (v3.2.1 — XLSX ile tutarlılık) ───
    rule = MaturityRule("Merge Hücreler", "Hücre", "Birleştirilmiş hücreler veri kaybına neden olabilir (processor otomatik çözümler)")
    total_merged = 0
    for sheet_idx in range(wb.nsheets):
        sheet = wb.sheet_by_index(sheet_idx)
        try:
            total_merged += len(sheet.merged_cells)
        except Exception:
            pass
    
    if total_merged > 20:
        rule.score = 50
        rule.status = "warning"
        rule.detail = f"{total_merged} birleştirilmiş hücre aralığı var. Çoğu otomatik çözümlenir ama karmaşık merge yapıları sorun çıkarabilir."
    elif total_merged > 5:
        rule.score = 75
        rule.status = "warning"
        rule.detail = f"{total_merged} birleştirilmiş hücre var. Otomatik çözümleme uygulanır."
    elif total_merged > 0:
        rule.score = 90
        rule.detail = f"{total_merged} birleştirilmiş hücre var (otomatik çözümlenir)."
    else:
        rule.score = 100
        rule.detail = "Birleştirilmiş hücre yok."
    rules.append(rule)
    
    return _build_result(rules, file_name, "XLS", total_sheets)


def analyze_txt(file_obj: BinaryIO, file_name: str) -> Dict[str, Any]:
    """TXT dosyasını analiz eder — basit kontroller"""
    rules: List[MaturityRule] = []
    
    file_obj.seek(0)
    content = file_obj.read()
    try:
        text = content.decode("utf-8")
    except UnicodeDecodeError:
        logger.debug("[MaturityAnalyzer] UTF-8 decode hatası, latin-1 fallback", exc_info=True)
        text = content.decode("latin-1", errors="replace")
    
    lines = text.split("\n")
    non_empty_lines = [l for l in lines if l.strip()]
    
    # Yapı kontrolü
    rule = MaturityRule("Yapısal Bölümleme", "Yapı", "Başlıklar veya bölüm ayraçları olmalı")
    heading_like = [l for l in non_empty_lines if len(l.strip()) < 60 and l.strip() and 
                    (l.strip().isupper() or re.match(r'^\d+[\.\)]\s', l.strip()) or l.strip().startswith('#'))]
    if heading_like:
        rule.score = 100
        rule.detail = f"{len(heading_like)} başlık/bölüm tespit edildi."
    else:
        rule.score = 50
        rule.status = "warning"
        rule.detail = "Yapısal bölümleme tespit edilemedi."
    rules.append(rule)
    
    # Encoding
    rule = MaturityRule("Encoding", "Encoding", "UTF-8 encoding olmalı")
    try:
        content.decode("utf-8")
        rule.score = 100
        rule.detail = "UTF-8 encoding sorunsuz."
    except UnicodeDecodeError:
        logger.debug("[MaturityAnalyzer] Encoding kontrolü - UTF-8 değil", exc_info=True)
        rule.score = 50
        rule.status = "warning"
        rule.detail = "UTF-8 encoding değil. Karakter sorunları olabilir."
    rules.append(rule)
    
    # İçerik yoğunluğu
    rule = MaturityRule("İçerik Yoğunluğu", "Metin", "Yeterli metin içeriği olmalı")
    if len(text.strip()) < 100:
        rule.score = 20
        rule.status = "fail"
        rule.detail = "Çok az içerik var."
    elif len(text.strip()) < 500:
        rule.score = 60
        rule.status = "warning"
        rule.detail = "İçerik düşük."
    else:
        rule.score = 100
        rule.detail = f"{len(non_empty_lines)} satır, {len(text)} karakter."
    rules.append(rule)
    
    return _build_result(rules, file_name, "TXT", len(non_empty_lines))


def analyze_pptx(file_obj: BinaryIO, file_name: str) -> Dict[str, Any]:
    """PPTX dosyasını analiz eder"""
    rules: List[MaturityRule] = []
    
    try:
        from pptx import Presentation
        file_obj.seek(0)
        prs = Presentation(file_obj)
        
        total_slides = len(prs.slides)
        total_text_chars = 0
        total_notes_chars = 0
        total_images = 0
        total_tables = 0
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    total_text_chars += len(shape.text_frame.text)
                # Görsel sayısı
                if shape.shape_type is not None and shape.shape_type == 13:  # Picture
                    total_images += 1
                # Tablo sayısı
                if shape.has_table:
                    total_tables += 1
            # Speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text or ""
                total_notes_chars += len(notes_text.strip())
        
        # ─── KURAL 1: Metin yoğunluğu ───
        rule = MaturityRule("Metin İçeriği", "Metin", "Slaytlarda yeterli metin olmalı")
        chars_per_slide = total_text_chars / max(total_slides, 1)
        if chars_per_slide < 50:
            rule.score = 30
            rule.status = "fail"
            rule.detail = f"Slayt başına ortalama {chars_per_slide:.0f} karakter. Çok az metin, muhtemelen görsel ağırlıklı."
        elif chars_per_slide < 150:
            rule.score = 65
            rule.status = "warning"
            rule.detail = f"Slayt başına {chars_per_slide:.0f} karakter."
        else:
            rule.score = 100
            rule.detail = f"Slayt başına {chars_per_slide:.0f} karakter. Yeterli."
        rules.append(rule)
        
        # ─── KURAL 2: Yapı kontrolü ───
        rule = MaturityRule("Slayt Başlıkları", "Yapı", "Her slaytın bir başlığı olmalı")
        titled_slides = sum(1 for s in prs.slides if s.shapes.title and s.shapes.title.text.strip())
        ratio = titled_slides / max(total_slides, 1)
        if ratio >= 0.8:
            rule.score = 100
            rule.detail = f"{titled_slides}/{total_slides} slaytın başlığı var."
        elif ratio >= 0.5:
            rule.score = 65
            rule.status = "warning"
            rule.detail = f"Sadece {titled_slides}/{total_slides} slaytın başlığı var."
        else:
            rule.score = 35
            rule.status = "fail"
            rule.detail = f"Çoğu slaytın başlığı yok ({titled_slides}/{total_slides})."
        rules.append(rule)
        
        # ─── KURAL 3: Speaker Notes ─── (v3.3.0)
        rule = MaturityRule("Speaker Notes", "İçerik", "Slayt notları önemli bağlam sağlar, RAG kalitesini artırır")
        if total_notes_chars > 100:
            rule.score = 100
            rule.detail = f"Toplam {total_notes_chars} karakter speaker notes var. RAG için zengin bağlam."
        elif total_notes_chars > 0:
            rule.score = 75
            rule.detail = f"Az miktarda speaker notes ({total_notes_chars} karakter). Daha fazla not eklenmesi önerilir."
        else:
            rule.score = 50
            rule.status = "warning"
            rule.detail = "Speaker notes boş. Slayt notlarına açıklayıcı bilgiler eklemek RAG kalitesini artırır."
        rules.append(rule)
        
        # ─── KURAL 4: Görsel/Metin Oranı ─── (v3.3.0)
        rule = MaturityRule("Görsel/Metin Oranı", "İçerik", "Görseller RAG tarafından işlenebilir ama metin daha etkili")
        if total_images > 0 and total_text_chars > 0:
            img_per_slide = total_images / max(total_slides, 1)
            if img_per_slide > 2 and chars_per_slide < 100:
                rule.score = 40
                rule.status = "fail"
                rule.detail = f"{total_images} görsel, slayt başına {chars_per_slide:.0f} karakter. Çok görsel ağırlıklı — metin bilgisi yetersiz."
            elif img_per_slide > 1.5:
                rule.score = 70
                rule.status = "warning"
                rule.detail = f"{total_images} görsel var. Görsellerdeki bilgileri alt metin olarak da ekleyin."
            else:
                rule.score = 100
                rule.detail = f"Görsel/metin dengesi uygun ({total_images} görsel)."
        else:
            rule.score = 100
            rule.detail = "Sorunsuz."
        rules.append(rule)
        
        # ─── KURAL 5: Tablo İçeriği ─── (v3.3.0)
        rule = MaturityRule("Tablo İçeriği", "Yapı", "Tablolar otomatik metin olarak çıkarılır")
        if total_tables > 0:
            rule.score = 90
            rule.detail = f"{total_tables} tablo tespit edildi. Tablolar metin olarak çıkarılacaktır."
        else:
            rule.score = 100
            rule.detail = "Tablo içeriği tespit edilmedi."
        rules.append(rule)
        
        # ─── KURAL 6: Slayt Sayısı ─── (v3.3.0)
        rule = MaturityRule("Slayt Sayısı", "Performans", "Yeterli slayt sayısı olmalı")
        if total_slides < 3:
            rule.score = 40
            rule.status = "warning"
            rule.detail = f"Sadece {total_slides} slayt var. Çok az içerik — RAG faydasız olabilir."
        elif total_slides > 200:
            rule.score = 65
            rule.status = "warning"
            rule.detail = f"{total_slides} slayt var. Çok büyük sunum — işleme süresi uzun olabilir."
        else:
            rule.score = 100
            rule.detail = f"{total_slides} slayt. Normal boyut."
        rules.append(rule)
        
    except Exception as e:
        rule = MaturityRule("Dosya Okunabilirlik", "Yapı", "Dosya okunabilir olmalı")
        rule.score = 50
        rule.status = "warning"
        rule.detail = f"Analiz sırasında sorun: {str(e)[:100]}"
        rules.append(rule)
    
    return _build_result(rules, file_name, "PPTX", 0)


def analyze_file(file_obj: BinaryIO, file_name: str) -> Dict[str, Any]:
    """
    Ana analiz fonksiyonu — dosya türüne göre uygun analizi çağırır.
    v3.3.0 [B3]: MD5 hash bazlı cache (10 dk TTL) — aynı dosyanın tekrar analizi önlenir.
    """
    import hashlib
    import time
    
    # Cache kontrolü — dosyanın ilk 1MB'ının MD5 hash'i
    file_obj.seek(0)
    first_mb = file_obj.read(1024 * 1024)
    file_obj.seek(0)
    
    file_hash = hashlib.md5(first_mb).hexdigest()
    cache_key = f"{file_hash}_{file_name}"
    
    now = time.time()
    CACHE_TTL = 600  # 10 dakika
    
    # Cache'den kontrol
    if cache_key in _maturity_cache:
        cached_result, cached_at = _maturity_cache[cache_key]
        if now - cached_at < CACHE_TTL:
            logger.debug("[MaturityAnalyzer] Cache hit: %s (hash=%s)", file_name, file_hash[:8])
            return cached_result
        else:
            del _maturity_cache[cache_key]
    
    # Cache temizliği (TTL geçmiş girişleri sil)
    expired_keys = [k for k, (_, t) in _maturity_cache.items() if now - t >= CACHE_TTL]
    for k in expired_keys:
        del _maturity_cache[k]
    
    ext = file_name.rsplit('.', 1)[-1].lower() if '.' in file_name else ''
    
    analyzers = {
        'pdf': analyze_pdf,
        'docx': analyze_docx,
        'doc': analyze_docx,
        'xlsx': analyze_xlsx,
        'xls': analyze_xls,
        'pptx': analyze_pptx,
        'ppt': analyze_pptx,
        'txt': analyze_txt,
        'csv': analyze_txt,
    }
    
    analyzer = analyzers.get(ext)
    if not analyzer:
        return {
            "file_name": file_name,
            "file_type": ext.upper(),
            "total_score": 75,
            "categories": [],
            "violations": [],
            "detail_count": 0,
            "message": "Bu dosya türü için analiz desteklenmiyor."
        }
    
    try:
        result = analyzer(file_obj, file_name)
        # Sonucu cache'e ekle
        _maturity_cache[cache_key] = (result, now)
        return result
    except Exception:
        logger.error("[MaturityAnalyzer] Analiz hatası: %s", file_name, exc_info=True)
        return {
            "file_name": file_name,
            "file_type": ext.upper(),
            "total_score": 50,
            "categories": [],
            "violations": [],
            "detail_count": 0,
            "message": "Dosya analizi sırasında bir hata oluştu."
        }


def _build_result(rules: List[MaturityRule], file_name: str, file_type: str, detail_count: int) -> Dict[str, Any]:
    """Kural sonuçlarından nihai skor ve kategori raporu oluşturur"""
    # Kategorilere grupla
    categories_map = {}
    for r in rules:
        if r.category not in categories_map:
            categories_map[r.category] = []
        categories_map[r.category].append(r)
    
    categories = []
    for cat_name, cat_rules in categories_map.items():
        cat_score = sum(r.score * r.weight for r in cat_rules) / sum(r.weight for r in cat_rules)
        categories.append({
            "name": cat_name,
            "score": round(cat_score, 1),
            "rules": [r.to_dict() for r in cat_rules]
        })
    
    # Toplam skor — tüm kuralların ağırlıklı ortalaması
    total_weight = sum(r.weight for r in rules)
    total_score = sum(r.score * r.weight for r in rules) / total_weight if total_weight > 0 else 0
    
    # İhlaller (skor < 100 olanlar)
    violations = [r.to_dict() for r in rules if r.status in ("fail", "warning")]
    
    return {
        "file_name": file_name,
        "file_type": file_type,
        "total_score": round(total_score, 1),
        "categories": categories,
        "violations": violations,
        "detail_count": detail_count,
        "message": None
    }
