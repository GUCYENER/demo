"""
VYRA L1 Support API - Image Extractor Module
=============================================
DOCX, PDF ve PPTX dosyalarından görsel çıkarma.
Görseller DB'de BYTEA olarak saklanır, chunk referansı ile ilişkilendirilir.

Author: VYRA AI Team
Version: 1.0.0
"""

import io
from dataclasses import dataclass
from typing import List, Optional

from app.services.logging_service import log_system_event, log_error


@dataclass
class ExtractedImage:
    """Çıkarılan görsel verisi"""
    image_data: bytes
    image_format: str       # png, jpeg, gif, emf
    width: int = 0
    height: int = 0
    context_heading: str = ""
    context_chunk_index: int = 0
    alt_text: str = ""
    ocr_text: str = ""      # EasyOCR ile çıkarılan metin
    paragraph_index: int = -1    # Görselin orijinal dosyadaki paragraf sırası (0-indexed)
    page_y_position: float = -1  # PDF: sayfa içi Y koordinatı (mm)


# Singleton OCR reader (lazy load — model yüklemesi ~20sn sadece ilk kez)
_ocr_reader = None


class ImageExtractor:
    """
    DOCX/PDF/PPTX dosyalarından görsel çıkarma modülü.
    
    Kullanım:
        extractor = ImageExtractor()
        images = extractor.extract(file_content, ".docx")
        if images:
            image_ids = extractor.save_to_db(images, file_id, cursor)
    """

    # Desteklenen görsel formatları
    SUPPORTED_FORMATS = {"png", "jpeg", "jpg", "gif", "bmp", "tiff", "emf", "wmf"}
    # OCR desteklenen formatlar (EMF/WMF Pillow tarafından açılamayabilir)
    OCR_FORMATS = {"png", "jpeg", "jpg", "gif", "bmp", "tiff"}

    def extract(self, file_content: bytes, file_type: str) -> List[ExtractedImage]:
        """
        Dosya tipine göre görselleri çıkarır.
        
        Args:
            file_content: Dosya binary içeriği
            file_type: Dosya uzantısı (örn: ".docx", ".pdf")
        
        Returns:
            Çıkarılan görsellerin listesi
        """
        ext = file_type.lower().strip(".")
        
        try:
            if ext in ("docx", "doc"):
                images = self._extract_docx_images(file_content)
            elif ext == "pdf":
                images = self._extract_pdf_images(file_content)
            elif ext in ("pptx", "ppt"):
                images = self._extract_pptx_images(file_content)
            else:
                return []
            
            # Tüm görseller için OCR çalıştır
            if images:
                self._run_ocr_batch(images)
            
            return images
        except Exception as e:
            log_error(f"Görsel çıkarma hatası ({ext}): {e}", "image_extractor")
            return []

    # ─────────────────────────────────────────
    #  DOCX Görsel Çıkarma
    # ─────────────────────────────────────────

    def _extract_docx_images(self, file_content: bytes) -> List[ExtractedImage]:
        """python-docx ile DOCX'ten görsel çıkar — paragraf pozisyonu da kaydedilir"""
        from docx import Document

        doc = Document(io.BytesIO(file_content))
        images: List[ExtractedImage] = []
        current_heading = "Genel"
        chunk_index = 0

        for para_idx, para in enumerate(doc.paragraphs):
            style_name = para.style.name if para.style else ""

            # Heading takibi
            if style_name.startswith("Heading") and para.text.strip():
                current_heading = para.text.strip()
                chunk_index += 1

            # InlineShape (gömülü resimler) kontrolü
            for run in para.runs:
                inline_shapes = run._element.findall(
                    './/{http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing}inline'
                )
                for inline in inline_shapes:
                    blip = inline.find(
                        './/{http://schemas.openxmlformats.org/drawingml/2006/main}blip'
                    )
                    if blip is not None:
                        embed_id = blip.get(
                            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
                        )
                        if embed_id:
                            img = self._get_docx_image_from_rel(
                                doc, embed_id, current_heading, chunk_index, para_idx
                            )
                            if img:
                                images.append(img)

        log_system_event("INFO", f"DOCX'ten {len(images)} görsel çıkarıldı", "image_extractor")
        return images

    def _get_docx_image_from_rel(self, doc, embed_id: str, heading: str, chunk_idx: int, para_idx: int = -1) -> Optional[ExtractedImage]:
        """DOCX relationship'ten görsel verisini alır"""
        try:
            rel = doc.part.rels.get(embed_id)
            if rel is None:
                return None

            image_part = rel.target_part
            image_data = image_part.blob
            content_type = image_part.content_type or ""

            # Format belirle
            fmt = self._format_from_content_type(content_type)
            if not fmt:
                return None

            width, height = self._get_image_dimensions(image_data, fmt)

            # Çok küçük görselleri atla (ikon, decoration vb.)
            if width < 50 or height < 50:
                return None

            return ExtractedImage(
                image_data=image_data,
                image_format=fmt,
                width=width,
                height=height,
                context_heading=heading[:500],
                context_chunk_index=chunk_idx,
                alt_text=f"Görsel - {heading[:100]}",
                paragraph_index=para_idx
            )
        except Exception as e:
            log_error(f"DOCX görsel okuma hatası: {e}", "image_extractor")
            return None

    # ─────────────────────────────────────────
    #  PDF Görsel Çıkarma
    # ─────────────────────────────────────────

    def _extract_pdf_images(self, file_content: bytes) -> List[ExtractedImage]:
        """PyMuPDF (fitz) ile PDF'ten görsel çıkar — sayfa heading'leri ve Y pozisyonu kaydeder"""
        import fitz
        import re

        doc = fitz.open(stream=file_content, filetype="pdf")
        images: List[ExtractedImage] = []
        
        # Toplam metin satır sayacı (yapay paragraph_index olarak kullanılacak)
        global_line_counter = 0

        for page_num in range(len(doc)):
            page = doc[page_num]
            image_list = page.get_images(full=True)
            
            if not image_list:
                # Sayfadaki satır sayısını sayacıya ekle (görsel olmasa bile)
                page_text = page.get_text("text") or ""
                global_line_counter += len([l for l in page_text.split('\n') if l.strip()])
                continue
            
            # Sayfadaki metin bloklarını pozisyonlarıyla birlikte al
            page_text = page.get_text("text") or ""
            page_lines = [l.strip() for l in page_text.split('\n') if l.strip()]
            
            # Heading'leri satır pozisyonu (metin bloğu dict) ile topla
            headings_on_page = []
            try:
                text_blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)
                for block in text_blocks.get("blocks", []):
                    if block.get("type") != 0:  # Sadece text blokları
                        continue
                    for line_info in block.get("lines", []):
                        spans = line_info.get("spans", [])
                        if not spans:
                            continue
                        line_text = "".join(s.get("text", "") for s in spans).strip()
                        if not line_text or len(line_text) > 80 or len(line_text) < 3:
                            continue
                        # Font size ve bold kontrolü → heading olabilir
                        avg_size = sum(s.get("size", 11) for s in spans) / len(spans)
                        is_bold = any(s.get("flags", 0) & 2**4 for s in spans)
                        y_pos = line_info["bbox"][1]  # Top Y
                        
                        if avg_size >= 12 or is_bold:
                            if (line_text[0].isupper() and
                                not line_text.endswith('.') and
                                not line_text.endswith(',') and
                                re.match(r'^[\d\.\s]*[A-ZÇĞİÖŞÜa-zçğıöşü]', line_text)):
                                headings_on_page.append((y_pos, line_text))
            except Exception:
                # Fallback: basit heading tespiti
                for line in page_lines:
                    if (len(line) < 80 and line[0].isupper() and 
                        not line.endswith('.') and
                        re.match(r'^[\d\.\s]*[A-ZÇĞİÖŞÜa-zçğıöşü]', line)):
                        headings_on_page.append((0, line))
            
            # Görsellerin sayfa üzerindeki bounding box'larını al
            for img_idx, img_info in enumerate(image_list):
                xref = img_info[0]
                try:
                    base_image = doc.extract_image(xref)
                    if not base_image or not base_image.get("image"):
                        continue

                    image_data = base_image["image"]
                    fmt = base_image.get("ext", "png").lower()

                    if fmt not in self.SUPPORTED_FORMATS:
                        continue

                    # Çok küçük görselleri atla (ikon, decoration vb.)
                    width = base_image.get("width", 0)
                    height = base_image.get("height", 0)
                    if width < 50 or height < 50:
                        continue
                    
                    # Görselin sayfa üzerindeki Y pozisyonunu bul
                    y_pos = 0.0
                    try:
                        img_rects = page.get_image_rects(xref)
                        if img_rects:
                            y_pos = img_rects[0].y0  # İlk rect'in top Y
                    except Exception:
                        y_pos = float(img_idx * 100)  # Fallback: sıra bazlı

                    # Bu görselin hemen öncesindeki heading'i bul
                    best_heading = f"Sayfa {page_num + 1}"
                    for h_y, h_text in sorted(headings_on_page, key=lambda x: x[0], reverse=True):
                        if h_y < y_pos:
                            best_heading = h_text
                            break
                    
                    # Yapay paragraph_index: görselden önceki metin satır sayısı
                    lines_before_img = sum(1 for l in page_lines if l.strip())
                    if y_pos > 0:
                        # Y pozisyonuna göre yaklaşık satır hesaplama
                        page_height = page.rect.height
                        lines_before_img = int((y_pos / page_height) * len(page_lines)) if page_height > 0 else 0
                    
                    para_idx = global_line_counter + lines_before_img

                    images.append(ExtractedImage(
                        image_data=image_data,
                        image_format=fmt if fmt != "jpg" else "jpeg",
                        width=width,
                        height=height,
                        context_heading=best_heading,
                        context_chunk_index=page_num,
                        alt_text=f"PDF Sayfa {page_num + 1} - Görsel {img_idx + 1}",
                        paragraph_index=para_idx,
                        page_y_position=y_pos
                    ))
                except Exception as e:
                    log_error(f"PDF görsel çıkarma hatası (sayfa {page_num}): {e}", "image_extractor")

            # Bu sayfanın satır sayısını toplam sayaca ekle
            global_line_counter += len(page_lines)

        doc.close()
        log_system_event("INFO", f"PDF'ten {len(images)} görsel çıkarıldı", "image_extractor")
        return images

    # ─────────────────────────────────────────
    #  PPTX Görsel Çıkarma
    # ─────────────────────────────────────────

    def _extract_pptx_images(self, file_content: bytes) -> List[ExtractedImage]:
        """python-pptx ile PPTX'ten görsel çıkar"""
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(io.BytesIO(file_content))
        images: List[ExtractedImage] = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_title = ""
            # Slayt başlığını bul
            if slide.shapes.title:
                slide_title = slide.shapes.title.text or f"Slayt {slide_idx + 1}"

            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        image_data = image.blob
                        content_type = image.content_type or ""

                        fmt = self._format_from_content_type(content_type)
                        if not fmt:
                            continue

                        width, height = self._get_image_dimensions(image_data, fmt)

                        # Çok küçük görselleri atla
                        if width < 50 or height < 50:
                            continue

                        images.append(ExtractedImage(
                            image_data=image_data,
                            image_format=fmt,
                            width=width,
                            height=height,
                            context_heading=(slide_title or f"Slayt {slide_idx + 1}")[:500],
                            context_chunk_index=slide_idx,
                            alt_text=f"Slayt {slide_idx + 1} - {slide_title[:100]}"
                        ))
                    except Exception as e:
                        log_error(f"PPTX görsel çıkarma hatası (slayt {slide_idx}): {e}", "image_extractor")

        log_system_event("INFO", f"PPTX'ten {len(images)} görsel çıkarıldı", "image_extractor")
        return images

    # ─────────────────────────────────────────
    #  Yardımcı Fonksiyonlar
    # ─────────────────────────────────────────

    def _format_from_content_type(self, content_type: str) -> Optional[str]:
        """Content-Type'tan görsel formatını belirle"""
        ct_map = {
            "image/png": "png",
            "image/jpeg": "jpeg",
            "image/jpg": "jpeg",
            "image/gif": "gif",
            "image/bmp": "bmp",
            "image/tiff": "tiff",
            "image/x-emf": "emf",
            "image/x-wmf": "wmf",
        }
        return ct_map.get(content_type.lower(), None)

    def _get_image_dimensions(self, image_data: bytes, fmt: str) -> tuple:
        """Görsel boyutlarını Pillow ile hesapla"""
        try:
            from PIL import Image
            img = Image.open(io.BytesIO(image_data))
            return img.size  # (width, height)
        except Exception as e:
            import sys
            print(f"[ImageExtractor] Image boyut okuma hatası: {e}", file=sys.stderr)
            return (0, 0)

    # ─────────────────────────────────────────
    #  OCR — EasyOCR ile Metin Çıkarma
    # ─────────────────────────────────────────

    def _get_ocr_reader(self):
        """Lazy singleton EasyOCR reader (ilk çağrıda model yüklenir)"""
        global _ocr_reader
        if _ocr_reader is None:
            try:
                import easyocr
                _ocr_reader = easyocr.Reader(['tr', 'en'], gpu=False, verbose=False)
                log_system_event("INFO", "EasyOCR reader yüklendi (tr+en)", "image_extractor")
            except ImportError:
                log_error("EasyOCR yüklü değil. pip install easyocr", "image_extractor")
                return None
            except Exception as e:
                log_error(f"EasyOCR yükleme hatası: {e}", "image_extractor")
                return None
        return _ocr_reader

    def _run_ocr_single(self, image_data: bytes, fmt: str) -> str:
        """
        Tek görsel için OCR çalıştır.
        
        Returns:
            Çıkarılan metin (boş string hata durumunda)
        """
        if fmt not in self.OCR_FORMATS:
            return ""
        
        reader = self._get_ocr_reader()
        if reader is None:
            return ""
        
        try:
            import numpy as np
            from PIL import Image
            
            img = Image.open(io.BytesIO(image_data))
            # RGB'ye çevir (RGBA, P vb. modlar EasyOCR'da sorun yaratabilir)
            if img.mode != 'RGB':
                img = img.convert('RGB')
            
            img_array = np.array(img)
            results = reader.readtext(img_array, detail=0, paragraph=True)
            
            text = "\n".join(results).strip()
            return text
        except Exception as e:
            log_error(f"OCR hatası: {e}", "image_extractor")
            return ""

    def _run_ocr_batch(self, images: List[ExtractedImage]):
        """
        Tüm görseller için OCR çalıştır ve ocr_text alanını doldur.
        ThreadPoolExecutor ile paralel çalışır (max 4 thread).
        """
        from concurrent.futures import ThreadPoolExecutor, as_completed
        
        if not images:
            return
        
        # İlk çağrıda model yüklenmesini sağla (thread-safe init)
        if self._get_ocr_reader() is None:
            log_system_event("WARNING", "OCR reader yüklenemedi, batch atlanıyor", "image_extractor")
            return
        
        max_workers = min(4, len(images))
        
        def _ocr_task(img):
            try:
                return self._run_ocr_single(img.image_data, img.image_format)
            except Exception as e:
                log_error(f"OCR task hatası ({img.image_format}): {e}", "image_extractor")
                return ""
        
        ocr_count = 0
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            future_to_img = {executor.submit(_ocr_task, img): img for img in images}
            for future in as_completed(future_to_img):
                img = future_to_img[future]
                try:
                    text = future.result()
                    if text:
                        img.ocr_text = text
                        ocr_count += 1
                except Exception as e:
                    log_error(f"OCR future hatası ({img.image_format}): {e}", "image_extractor")
        
        if ocr_count > 0:
            log_system_event("INFO", f"{ocr_count}/{len(images)} görselden OCR metin çıkarıldı (paralel)", "image_extractor")
        elif len(images) > 0:
            log_system_event("WARNING", f"0/{len(images)} görselden OCR metin çıkarılamadı", "image_extractor")

    def save_to_db(self, images: List[ExtractedImage], file_id: int, cursor) -> List[int]:
        """
        Görselleri document_images tablosuna kaydet.
        
        Args:
            images: Çıkarılan görseller
            file_id: Dosya ID
            cursor: DB cursor (dışarıdan geçirilir, commit yapılmaz)
        
        Returns:
            Eklenen görsel ID'lerinin listesi
        """
        if not images:
            return []

        image_ids = []
        for idx, img in enumerate(images):
            try:
                cursor.execute(
                    """
                    INSERT INTO document_images 
                    (file_id, image_index, image_data, image_format, 
                     width_px, height_px, file_size_bytes, 
                     context_heading, context_chunk_index, alt_text, ocr_text)
                    VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
                    RETURNING id
                    """,
                    (
                        file_id, idx, img.image_data, img.image_format,
                        img.width, img.height, len(img.image_data),
                        img.context_heading, img.context_chunk_index, img.alt_text,
                        img.ocr_text
                    )
                )
                row = cursor.fetchone()
                if row:
                    image_ids.append(row["id"])
            except Exception as e:
                log_error(f"Görsel DB kayıt hatası (index {idx}): {e}", "image_extractor")

        log_system_event("INFO", f"Dosya {file_id} için {len(image_ids)} görsel kaydedildi", "image_extractor")
        return image_ids
